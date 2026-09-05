# Minimal runnable fixtures for the snippets called out in review: one file per format.
# Each script is self-contained, writes only scratch files into the current directory,
# and exits non-zero on failed assertions. Run from any scratch directory:
#   python pdf_fixture.py   (deps: reportlab, pypdf, pymupdf)
#   python pptx_fixture.py  (deps: python-pptx)
#   python xlsx_fixture.py  (deps: openpyxl)
#   python docx_fixture.py  (deps: python-docx, pymupdf, soffice on PATH)
import base64
import copy
import sys
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import BubbleChartData, ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt

failures = []


def check(name, cond, extra=""):
    print(("PASS " if cond else "FAIL ") + name + ((" :: " + str(extra)) if not cond and extra else ""))
    if not cond:
        failures.append(name)


# ---- build the deck ------------------------------------------------------------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tf = box.text_frame
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "old wording"
r2 = p.add_run()
r2.text = " linked part"
r2.font.italic = True
r2.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
r2.hyperlink.address = "https://example.com/docs"

table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(6), Inches(1))
table = table_shape.table
cell = table.cell(0, 1)
ctf = cell.text_frame
cp = ctf.paragraphs[0]
cr = cp.add_run()
cr.text = "old cell text"
cr.font.bold = True
cr.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)

chart_data = ChartData()
chart_data.categories = ["EU", "US"]
chart_data.add_series("Units", (120, 80))
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7.2), Inches(2.5), Inches(2), Inches(2),
    chart_data,
).chart
chart.has_title = True
chart.chart_title.text_frame.text = "Units by region"
chart.category_axis.has_title = True
chart.category_axis.axis_title.text_frame.text = "Region"
chart.value_axis.has_title = True
chart.value_axis.axis_title.text_frame.text = "Units sold"

xy_data = XyChartData()
xy_series = xy_data.add_series("Trend")
xy_series.add_data_point(1, 2)
xy_series.add_data_point(3, 4)
xy_chart = slide.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER, Inches(7.2), Inches(4.7), Inches(2), Inches(2), xy_data,
).chart
xy_chart.has_title = True
xy_chart.chart_title.text_frame.text = "XY trend"

bubble_data = BubbleChartData()
bubble_series = bubble_data.add_series("Risk")
bubble_series.add_data_point(3, 4, 5)
bubble_chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE, Inches(5), Inches(4.7), Inches(2), Inches(2), bubble_data,
).chart
bubble_chart.has_title = True
bubble_chart.chart_title.text_frame.text = "Bubble risk"
slide.notes_slide.notes_text_frame.text = "Speaker note: explain the regional split."

prs.save("input.pptx")

# ---- create.md skeleton: explicit placeholder selection and populated table ----
def placeholder_of_type(slide, *types):
    matches = [
        ph for ph in slide.placeholders
        if ph.placeholder_format.type in types
    ]
    if len(matches) != 1:
        available = [
            f"{ph.name} ({ph.placeholder_format.type})"
            for ph in slide.placeholders
        ]
        raise ValueError(
            f"expected exactly one placeholder of {types}, found {len(matches)}; "
            f"available placeholders: {available or 'none'}"
        )
    return matches[0]


selector_prs = Presentation()
no_object_slide = selector_prs.slides.add_slide(selector_prs.slide_layouts[6])
try:
    placeholder_of_type(no_object_slide, PP_PLACEHOLDER.OBJECT)
    missing_placeholder_message = None
except ValueError as exc:
    missing_placeholder_message = str(exc)
check(
    "placeholder selector explicitly rejects zero matches",
    missing_placeholder_message is not None
    and "found 0" in missing_placeholder_message
    and "available placeholders: none" in missing_placeholder_message,
    missing_placeholder_message,
)

two_object_slide = selector_prs.slides.add_slide(selector_prs.slide_layouts[3])
try:
    placeholder_of_type(two_object_slide, PP_PLACEHOLDER.OBJECT)
    ambiguous_placeholder_message = None
except ValueError as exc:
    ambiguous_placeholder_message = str(exc)
check(
    "placeholder selector explicitly rejects multiple matches",
    ambiguous_placeholder_message is not None
    and "found 2" in ambiguous_placeholder_message
    and "Content Placeholder 2" in ambiguous_placeholder_message
    and "Content Placeholder 3" in ambiguous_placeholder_message,
    ambiguous_placeholder_message,
)

skeleton_prs = Presentation()
skeleton_slide = skeleton_prs.slides.add_slide(skeleton_prs.slide_layouts[5])
skeleton_slide.shapes.title.text = "Regional service health"
skeleton_headers = ["Region", "Error rate", "P99 latency"]
skeleton_body = [
    ["Americas", "0.08%", "182 ms"],
    ["Europe", "0.05%", "164 ms"],
    ["Asia Pacific", "0.11%", "213 ms"],
]
skeleton_table = skeleton_slide.shapes.add_table(
    1 + len(skeleton_body), len(skeleton_headers),
    Inches(0.5), Inches(1.5), Inches(9), Inches(3.5),
).table
for column_index, text in enumerate(skeleton_headers):
    skeleton_table.cell(0, column_index).text = text
for row_index, row in enumerate(skeleton_body, start=1):
    for column_index, text in enumerate(row):
        skeleton_table.cell(row_index, column_index).text = text
skeleton_prs.save("create-skeleton.pptx")
skeleton_reopened = Presentation("create-skeleton.pptx")
skeleton_reopened_table = next(
    shape.table for shape in skeleton_reopened.slides[0].shapes if shape.has_table
)
skeleton_values = [
    [cell.text.strip() for cell in row.cells]
    for row in skeleton_reopened_table.rows
]
check(
    "creation skeleton table has no unexpected empty cells",
    skeleton_values == [skeleton_headers, *skeleton_body]
    and all(text for row in skeleton_values for text in row),
    skeleton_values,
)

# ---- analyze.md bounded package check rejects archive bombs before expansion ---
MAX_ARCHIVE_BYTES = 200 * 1024 * 1024
MAX_MEMBERS = 10_000
MAX_XML_PART = 20 * 1024 * 1024
MAX_ENTRY = 100 * 1024 * 1024
MAX_TOTAL_UNCOMPRESSED = 500 * 1024 * 1024
MAX_COMPRESSION_RATIO = 200
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
SAFE_XML = etree.XMLParser(
    load_dtd=False,
    resolve_entities=False,
    no_network=True,
    huge_tree=False,
    recover=False,
)


def require(condition, message):
    if not condition:
        raise ValueError(message)


def validate_pptx_package(source):
    source.seek(0, 2)
    compressed_size = source.tell()
    source.seek(0)
    require(compressed_size <= MAX_ARCHIVE_BYTES,
            "compressed PPTX file size above limit")
    with zipfile.ZipFile(source) as archive:
        infos = archive.infolist()
        require(len(infos) <= MAX_MEMBERS, "archive member count above limit")
        names = {info.filename for info in infos}
        require(len(names) == len(infos), "duplicate archive member names are unsafe")
        require("[Content_Types].xml" in names and "ppt/presentation.xml" in names,
                "missing required OPC members")
        require(sum(info.file_size for info in infos) <= MAX_TOTAL_UNCOMPRESSED,
                "declared total uncompressed size above limit")

        content_types_info = next(
            info for info in infos if info.filename == "[Content_Types].xml"
        )
        require(content_types_info.file_size <= MAX_XML_PART,
                "oversized XML part: [Content_Types].xml")
        require(
            content_types_info.file_size / max(content_types_info.compress_size, 1)
            <= MAX_COMPRESSION_RATIO,
            "suspicious compression ratio: [Content_Types].xml",
        )
        content_type_chunks = []
        content_type_size = 0
        with archive.open(content_types_info) as stream:
            while chunk := stream.read(64 * 1024):
                content_type_size += len(chunk)
                require(content_type_size <= MAX_XML_PART,
                        "part exceeded read limit: [Content_Types].xml")
                content_type_chunks.append(chunk)
        require(content_type_size == content_types_info.file_size,
                "size mismatch: [Content_Types].xml")
        content_types_blob = b"".join(content_type_chunks)
        content_types_root = etree.fromstring(content_types_blob, parser=SAFE_XML)
        require(content_types_root.tag == f"{{{CONTENT_TYPES_NS}}}Types",
                "invalid content-types root")
        default_types = {}
        override_types = {}
        for declaration in content_types_root:
            if declaration.tag == f"{{{CONTENT_TYPES_NS}}}Default":
                key = (declaration.get("Extension") or "").casefold()
                target = default_types
            elif declaration.tag == f"{{{CONTENT_TYPES_NS}}}Override":
                part_name = declaration.get("PartName") or ""
                require(part_name.startswith("/"), "invalid content-type part name")
                key = part_name[1:].casefold()
                target = override_types
            else:
                continue
            content_type = declaration.get("ContentType") or ""
            require(key and content_type and key not in target,
                    "invalid or duplicate content-type declaration")
            target[key] = content_type.partition(";")[0].strip().casefold()

        actual_total = 0
        for info in infos:
            require(info.file_size <= MAX_ENTRY, f"oversized part: {info.filename}")
            ratio = info.file_size / max(info.compress_size, 1)
            require(ratio <= MAX_COMPRESSION_RATIO,
                    f"suspicious compression ratio: {info.filename}")
            extension = info.filename.rpartition(".")[2].casefold()
            content_type = override_types.get(
                info.filename.casefold(), default_types.get(extension, "")
            )
            is_xml = (
                info.filename.casefold().endswith((".xml", ".rels"))
                or content_type in {"application/xml", "text/xml"}
                or content_type.endswith("+xml")
            )
            if is_xml:
                require(info.file_size <= MAX_XML_PART, f"oversized XML part: {info.filename}")
            chunks = []
            actual_size = 0
            if info.filename == "[Content_Types].xml":
                chunks = [content_types_blob]
                actual_size = len(content_types_blob)
                actual_total += actual_size
            else:
                with archive.open(info) as stream:
                    while chunk := stream.read(64 * 1024):
                        actual_size += len(chunk)
                        actual_total += len(chunk)
                        require(actual_size <= MAX_ENTRY,
                                f"part exceeded read limit: {info.filename}")
                        require(actual_total <= MAX_TOTAL_UNCOMPRESSED,
                                "archive exceeded total read limit")
                        if is_xml:
                            chunks.append(chunk)
            require(actual_total <= MAX_TOTAL_UNCOMPRESSED,
                    "archive exceeded total read limit")
            require(actual_size == info.file_size, f"size mismatch: {info.filename}")
            if is_xml:
                etree.fromstring(b"".join(chunks), parser=SAFE_XML)
    source.seek(0)


def open_validated_presentation(path):
    with Path(path).open("rb") as source:
        validate_pptx_package(source)
        source.seek(0)
        return Presentation(source)


try:
    with Path("input.pptx").open("rb") as source:
        validate_pptx_package(source)
    healthy_pptx_passed = True
except Exception:
    healthy_pptx_passed = False
check("bounded PPTX check accepts an ordinary deck", healthy_pptx_passed)

Path("oversized-before-open.pptx").write_bytes(b"not a ZIP package")
original_archive_limit = MAX_ARCHIVE_BYTES
MAX_ARCHIVE_BYTES = 0
try:
    with Path("oversized-before-open.pptx").open("rb") as source:
        validate_pptx_package(source)
    compressed_size_rejected_before_open = False
except ValueError as exc:
    compressed_size_rejected_before_open = (
        str(exc) == "compressed PPTX file size above limit"
    )
finally:
    MAX_ARCHIVE_BYTES = original_archive_limit
check(
    "compressed PPTX size is bounded before ZipFile opens the package",
    compressed_size_rejected_before_open,
)

with zipfile.ZipFile("too-many-members.pptx", "w", zipfile.ZIP_STORED) as archive:
    for member_index in range(MAX_MEMBERS + 1):
        archive.writestr(f"zero-{member_index:05d}.bin", b"")
try:
    with Path("too-many-members.pptx").open("rb") as source:
        validate_pptx_package(source)
    many_members_rejected = False
except ValueError as exc:
    many_members_rejected = str(exc) == "archive member count above limit"
check(
    "member-count gate rejects 10,001 zero-byte members before traversal",
    many_members_rejected,
)

with zipfile.ZipFile("input.pptx") as source_archive:
    compressed_bomb_payload = {
        info.filename: source_archive.read(info) for info in source_archive.infolist()
    }
compressed_bomb_member = next(
    name for name in compressed_bomb_payload
    if not name.casefold().endswith((".xml", ".rels"))
)
compressed_bomb_payload[compressed_bomb_member] = b"x" * 2_000_000
with zipfile.ZipFile("compressed-bomb.pptx", "w", zipfile.ZIP_DEFLATED) as archive:
    for member_name, member_data in compressed_bomb_payload.items():
        archive.writestr(member_name, member_data)
try:
    with Path("compressed-bomb.pptx").open("rb") as source:
        validate_pptx_package(source)
    pptx_bomb_rejected = False
except ValueError as error:
    pptx_bomb_rejected = (
        str(error) == f"suspicious compression ratio: {compressed_bomb_member}"
    )
check("valid-manifest PPTX compression bomb reaches the ratio gate",
      pptx_bomb_rejected, compressed_bomb_member)

with zipfile.ZipFile("input.pptx") as source_archive:
    uppercase_xml_payload = {
        info.filename: source_archive.read(info) for info in source_archive.infolist()
    }
uppercase_slide = uppercase_xml_payload.pop("ppt/slides/slide1.xml")
uppercase_xml_limit = 64 * 1024
uppercase_xml_payload["ppt/slides/slide1.XML"] = uppercase_slide.replace(
    b"</p:sld>", b" " * uppercase_xml_limit + b"</p:sld>", 1,
)
uppercase_xml_payload["[Content_Types].xml"] = uppercase_xml_payload[
    "[Content_Types].xml"
].replace(b"/ppt/slides/slide1.xml", b"/ppt/slides/slide1.XML")
uppercase_xml_payload["ppt/_rels/presentation.xml.rels"] = uppercase_xml_payload[
    "ppt/_rels/presentation.xml.rels"
].replace(b"slides/slide1.xml", b"slides/slide1.XML")
with zipfile.ZipFile("uppercase-xml-part.pptx", "w", zipfile.ZIP_STORED) as archive:
    for member_name, member_data in uppercase_xml_payload.items():
        archive.writestr(member_name, member_data)
uppercase_route_parses = len(Presentation("uppercase-xml-part.pptx").slides) == 1
original_xml_limit = MAX_XML_PART
MAX_XML_PART = uppercase_xml_limit
try:
    with Path("uppercase-xml-part.pptx").open("rb") as source:
        validate_pptx_package(source)
    uppercase_oversized_xml_rejected = False
except ValueError as error:
    uppercase_oversized_xml_rejected = (
        str(error) == "oversized XML part: ppt/slides/slide1.XML"
    )
finally:
    MAX_XML_PART = original_xml_limit
check("uppercase XML part names cannot bypass the XML size and parser gate",
      uppercase_route_parses and uppercase_oversized_xml_rejected)

with zipfile.ZipFile("input.pptx") as source_archive:
    typed_xml_payload = {
        info.filename: source_archive.read(info) for info in source_archive.infolist()
    }
typed_slide = typed_xml_payload.pop("ppt/slides/slide1.xml")
typed_xml_payload["ppt/slides/Slide1.DaT"] = typed_slide.replace(
    b"</p:sld>", b" " * uppercase_xml_limit + b"</p:sld>", 1,
)
typed_xml_payload["[Content_Types].xml"] = typed_xml_payload[
    "[Content_Types].xml"
].replace(b"/ppt/slides/slide1.xml", b"/PPT/SLIDES/SLIDE1.DAT")
typed_xml_payload["ppt/_rels/presentation.xml.rels"] = typed_xml_payload[
    "ppt/_rels/presentation.xml.rels"
].replace(b"slides/slide1.xml", b"slides/Slide1.DaT")
with zipfile.ZipFile("content-typed-xml-part.pptx", "w", zipfile.ZIP_STORED) as archive:
    for member_name, member_data in typed_xml_payload.items():
        archive.writestr(member_name, member_data)
typed_xml_route_parses = len(Presentation("content-typed-xml-part.pptx").slides) == 1
original_xml_limit = MAX_XML_PART
MAX_XML_PART = uppercase_xml_limit
try:
    with Path("content-typed-xml-part.pptx").open("rb") as source:
        validate_pptx_package(source)
    typed_oversized_xml_rejected = False
except ValueError as error:
    typed_oversized_xml_rejected = (
        str(error) == "oversized XML part: ppt/slides/Slide1.DaT"
    )
finally:
    MAX_XML_PART = original_xml_limit
check("case-insensitive XML content types cannot bypass bounds with an arbitrary extension",
      typed_xml_route_parses and typed_oversized_xml_rejected)

mixed_case_xml_rejections = []
for mixed_case_name in ("custom.XmL", "custom.ReLs"):
    with zipfile.ZipFile("input.pptx") as source_archive:
        mixed_case_payload = {
            info.filename: source_archive.read(info) for info in source_archive.infolist()
        }
    mixed_case_payload[mixed_case_name] = (
        b"<root>" + b" " * uppercase_xml_limit + b"</root>"
    )
    mixed_case_declaration = (
        f'<Override PartName="/{mixed_case_name}" '
        'ContentType="application/octet-stream"/>'
    ).encode()
    mixed_case_payload["[Content_Types].xml"] = mixed_case_payload[
        "[Content_Types].xml"
    ].replace(b"</Types>", mixed_case_declaration + b"</Types>", 1)
    mixed_case_path = "mixed-case-" + mixed_case_name.replace(".", "-") + ".pptx"
    with zipfile.ZipFile(mixed_case_path, "w", zipfile.ZIP_STORED) as archive:
        for member_name, member_data in mixed_case_payload.items():
            archive.writestr(member_name, member_data)
    original_xml_limit = MAX_XML_PART
    MAX_XML_PART = uppercase_xml_limit
    try:
        with Path(mixed_case_path).open("rb") as source:
            validate_pptx_package(source)
    except ValueError as error:
        if str(error) == f"oversized XML part: {mixed_case_name}":
            mixed_case_xml_rejections.append(mixed_case_name)
    finally:
        MAX_XML_PART = original_xml_limit
check("mixed-case XML and relationship suffixes keep XML bounds",
      mixed_case_xml_rejections == ["custom.XmL", "custom.ReLs"],
      mixed_case_xml_rejections)

with zipfile.ZipFile("input.pptx") as source_archive:
    default_xml_payload = {
        info.filename: source_archive.read(info) for info in source_archive.infolist()
    }
default_xml_payload["custom.payload"] = (
    b"<root>" + b" " * uppercase_xml_limit + b"</root>"
)
default_xml_declaration = (
    b'<Default Extension="payload" ContentType="application/custom+xml"/>'
)
default_xml_payload["[Content_Types].xml"] = default_xml_payload[
    "[Content_Types].xml"
].replace(b"</Types>", default_xml_declaration + b"</Types>", 1)
with zipfile.ZipFile("default-content-type-xml.pptx", "w", zipfile.ZIP_STORED) as archive:
    for member_name, member_data in default_xml_payload.items():
        archive.writestr(member_name, member_data)
original_xml_limit = MAX_XML_PART
MAX_XML_PART = uppercase_xml_limit
try:
    with Path("default-content-type-xml.pptx").open("rb") as source:
        validate_pptx_package(source)
    default_typed_xml_rejected = False
except ValueError as error:
    default_typed_xml_rejected = str(error) == "oversized XML part: custom.payload"
finally:
    MAX_XML_PART = original_xml_limit
check("default +xml content types apply XML bounds to arbitrary extensions",
      default_typed_xml_rejected)

duplicate_override_payload = dict(typed_xml_payload)
duplicate_override = (
    b'<Override PartName="/ppt/slides/slide1.dat" '
    b'ContentType="application/octet-stream"/>'
)
duplicate_override_payload["[Content_Types].xml"] = duplicate_override_payload[
    "[Content_Types].xml"
].replace(b"</Types>", duplicate_override + b"</Types>", 1)
with zipfile.ZipFile("duplicate-case-override.pptx", "w", zipfile.ZIP_STORED) as archive:
    for member_name, member_data in duplicate_override_payload.items():
        archive.writestr(member_name, member_data)
try:
    with Path("duplicate-case-override.pptx").open("rb") as source:
        validate_pptx_package(source)
    duplicate_case_override_rejected = False
except ValueError as error:
    duplicate_case_override_rejected = (
        str(error) == "invalid or duplicate content-type declaration"
    )
check("case-insensitive duplicate content-type overrides fail closed",
      duplicate_case_override_rejected)

presentation_loader_calls = []
real_presentation_loader = Presentation


def unexpected_presentation_loader(source):
    presentation_loader_calls.append(source)
    return real_presentation_loader(source)


Presentation = unexpected_presentation_loader
try:
    open_validated_presentation("compressed-bomb.pptx")
    bomb_rejected_before_presentation = False
except ValueError:
    bomb_rejected_before_presentation = not presentation_loader_calls
finally:
    Presentation = real_presentation_loader
check("content inventory rejects a package bomb before Presentation parses it",
      bomb_rejected_before_presentation, len(presentation_loader_calls))

validated_sources = []
parsed_sources = []
parsed_source_open_states = []
real_package_validator = validate_pptx_package
real_presentation_loader = Presentation


def tracking_package_validator(source):
    validated_sources.append(source)
    return real_package_validator(source)


def tracking_presentation_loader(source):
    parsed_sources.append(source)
    parsed_source_open_states.append(not source.closed)
    return real_presentation_loader(source)


validate_pptx_package = tracking_package_validator
Presentation = tracking_presentation_loader
try:
    same_handle_prs = open_validated_presentation("input.pptx")
finally:
    validate_pptx_package = real_package_validator
    Presentation = real_presentation_loader
check("validated presentation parses the exact same open handle before closing it",
      len(same_handle_prs.slides) == 1
      and len(validated_sources) == len(parsed_sources) == 1
      and validated_sources[0] is parsed_sources[0]
      and parsed_source_open_states == [True]
      and validated_sources[0].closed)

# ---- edit.md snippet: validated single-shape run replace keeps styling/link ----
edit_loader_calls = []
real_presentation_loader = Presentation


def unexpected_edit_loader(source):
    edit_loader_calls.append(source)
    return real_presentation_loader(source)


Presentation = unexpected_edit_loader
try:
    open_validated_presentation("compressed-bomb.pptx")
    edit_bomb_rejected_before_parse = False
except ValueError:
    edit_bomb_rejected_before_parse = not edit_loader_calls
finally:
    Presentation = real_presentation_loader
check("edit route rejects a package bomb before Presentation parses it",
      edit_bomb_rejected_before_parse, edit_loader_calls)

prs = open_validated_presentation("input.pptx")
old, new = "old wording", "new wording"
candidates = []
for i, s in enumerate(prs.slides):
    for shape in s.shapes:
        if shape.has_text_frame and old in shape.text_frame.text:
            candidates.append((i, shape.name, shape))
require(len(candidates) == 1, "expected exactly one text target")
_, _, target_shape = candidates[0]
tf = target_shape.text_frame
run_hits = [run for par in tf.paragraphs for run in par.runs if old in run.text]
require(len(run_hits) == 1, "target is duplicated or split across runs")
run_hits[0].text = run_hits[0].text.replace(old, new, 1)

edited_link = [run for par in tf.paragraphs for run in par.runs if run.hyperlink.address]
check("run replace keeps the other run's hyperlink", len(edited_link) == 1 and edited_link[0].hyperlink.address == "https://example.com/docs")
styled = [run for par in tf.paragraphs for run in par.runs if run.font.italic]
check("run replace keeps sibling run styling", len(styled) == 1 and styled[0].font.color.rgb == RGBColor(0xC0, 0x00, 0x00))
prs.save("edited.pptx")

# ---- edit.md snippet: table cell edited at run level ---------------------------
prs2 = Presentation("input.pptx")
old_cell, new_cell = "old cell text", "new cell text"
tbl = next(sh for sh in prs2.slides[0].shapes if sh.has_table).table
cell = tbl.cell(0, 1)
hits = [run for par in cell.text_frame.paragraphs for run in par.runs if old_cell in run.text]
require(len(hits) == 1, "target is duplicated or split across runs in this cell")
hits[0].text = hits[0].text.replace(old_cell, new_cell, 1)

prs2.save("cell-edited.pptx")
prs3 = Presentation("cell-edited.pptx")
cell3 = next(sh for sh in prs3.slides[0].shapes if sh.has_table).table.cell(0, 1)
after_runs = [run for par in cell3.text_frame.paragraphs for run in par.runs]
check("cell run edit keeps bold", any(r.font.bold for r in after_runs))
check("cell run edit keeps color", any(r.font.color and r.font.color.rgb == RGBColor(0x00, 0x70, 0xC0) for r in after_runs))
check("cell run edit changed the text", cell3.text_frame.text == "new cell text")

# the dangerous variant for contrast: assigning cell.text drops run properties
prs4 = Presentation("input.pptx")
tbl4 = next(sh for sh in prs4.slides[0].shapes if sh.has_table).table
tbl4.cell(0, 1).text = "new cell text"
prs4.save("cell-flattened.pptx")
prs5 = Presentation("cell-flattened.pptx")
flat_runs = [run for par in next(sh for sh in prs5.slides[0].shapes if sh.has_table).table.cell(0, 1).text_frame.paragraphs for run in par.runs]
check("cell.text assignment is proven lossy (negative control)", not any(r.font.bold for r in flat_runs))

# ---- analyze.md snippet: grouped-shape walker ----------------------------------
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

prs6 = Presentation()
slide6 = prs6.slides.add_slide(prs6.slide_layouts[5])
pic_holder = slide6.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
pic_holder.text_frame.text = "nested member"
pic_holder.text_frame.paragraphs[0].runs[0].font.name = "Grouped Face"

GRP = (
    '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    "<p:nvGrpSpPr><p:cNvPr id=\"901\" name=\"demo group\"/>"
    "<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>"
    "<p:grpSpPr><a:xfrm><a:off x=\"914400\" y=\"914400\"/>"
    '<a:ext cx="4572000" cy="914400"/>'
    '<a:chOff x="0" y="0"/><a:chExt cx="4572000" cy="914400"/></a:xfrm></p:grpSpPr>'
    "</p:grpSp>"
)


def iter_shapes(shapes):
    for shape in shapes:
        if shape_is_hidden(shape):
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


OOXML_TRUE = {"1", "true"}
OOXML_FALSE = {"0", "false"}


def ooxml_bool(element, attribute, default):
    value = element.get(attribute)
    if value is None:
        return default
    value = value.strip(" \t\r\n")
    if value in OOXML_TRUE:
        return True
    if value in OOXML_FALSE:
        return False
    raise ValueError(f"invalid OOXML boolean {attribute}={value!r}")


def shape_is_hidden(shape):
    properties = shape._element.find(".//" + qn("p:cNvPr"))
    return properties is not None and ooxml_bool(properties, "hidden", False)


def layer_text_content(shapes, source, *, inherited):
    records = []
    for shape in shapes:
        if shape_is_hidden(shape):
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            records.extend(layer_text_content(shape.shapes, source, inherited=inherited))
            continue
        if inherited and shape.is_placeholder:
            continue
        if shape.has_text_frame and shape.text_frame.text:
            records.append({"source": source, "shape": shape.name,
                            "text": shape.text_frame.text})
    return records


def slide_text_content(slide):
    records = []
    if ooxml_bool(slide._element, "showMasterSp", True):
        layout = slide.slide_layout
        if ooxml_bool(layout._element, "showMasterSp", True):
            records.extend(layer_text_content(
                layout.slide_master.shapes, "master", inherited=True,
            ))
        records.extend(layer_text_content(layout.shapes, "layout", inherited=True))
    records.extend(layer_text_content(slide.shapes, "slide", inherited=False))
    return records


def table_cells(table):
    return [[
        {
            "row": row_index,
            "column": column_index,
            "text": None if cell.is_spanned else cell.text,
            "is_merge_origin": cell.is_merge_origin,
            "is_spanned": cell.is_spanned,
            "span_width": cell.span_width,
            "span_height": cell.span_height,
        }
        for column_index, cell in enumerate(row.cells)
    ] for row_index, row in enumerate(table.rows)]


def picture_content(shape):
    if getattr(type(shape), "image", None) is None:
        return None
    try:
        image = shape.image
        return {
            "name": shape.name,
            "filename": image.filename,
            "extension": image.ext,
            "bytes": len(image.blob),
        }
    except (AttributeError, KeyError, OSError, ValueError):
        blips = shape._element.xpath(".//a:blip")
        relationship_id = (
            blips[0].get(qn("r:embed")) if len(blips) == 1 else None
        )
        return {
            "name": shape.name,
            "status": "unreadable",
            "relationship_id": relationship_id,
            "reason": "missing or invalid image relationship or payload",
        }


MAX_CHART_POINTS = 100_000
MAX_CATEGORY_LEVELS = 16
CATEGORY_SOURCE_NAMES = {"strRef", "numRef", "multiLvlStrRef", "strLit", "numLit"}


def cache_point_count(cache):
    point_counts = cache.findall(qn("c:ptCount"))
    if len(point_counts) != 1:
        return None
    try:
        point_count = int(point_counts[0].get("val"))
    except (TypeError, ValueError):
        return None
    return point_count if 0 <= point_count <= MAX_CHART_POINTS else None


def consume_point_budget(point_budget, count):
    if count > MAX_CHART_POINTS:
        return False
    if point_budget is None:
        return True
    if count > point_budget["remaining"]:
        return False
    point_budget["remaining"] -= count
    return True


def cached_numeric_points(
    source, *, fill_missing=False, include_count=False, point_budget=None
):
    if source is None:
        return None
    caches = source.xpath("./c:numRef/c:numCache | ./c:numLit")
    if len(caches) != 1 or (point_count := cache_point_count(caches[0])) is None:
        return None
    if not consume_point_budget(point_budget, point_count):
        return None
    cached = {}
    for point in caches[0].findall(qn("c:pt")):
        value = point.find(qn("c:v"))
        try:
            index = int(point.get("idx"))
        except (TypeError, ValueError):
            return None
        if value is None or not 0 <= index < point_count or index in cached:
            return None
        try:
            parsed = float(value.text) if value.text not in (None, "") else value.text
        except (TypeError, ValueError):
            parsed = value.text
        cached[index] = parsed
    if fill_missing:
        points = [(index, cached.get(index)) for index in range(point_count)]
    else:
        points = sorted(cached.items())
    return {"point_count": point_count, "points": points} if include_count else points


def cached_text_point_map(container, point_count):
    values = {}
    for point in container.findall(qn("c:pt")):
        try:
            index = int(point.get("idx"))
        except (TypeError, ValueError):
            return None
        value_nodes = point.findall(qn("c:v"))
        if (len(value_nodes) != 1 or not 0 <= index < point_count
                or index in values):
            return None
        values[index] = value_nodes[0].text or ""
    return values


def category_source(series):
    category = series._element.find(qn("c:cat"))
    if category is None:
        return None, None
    sources = [
        child for child in category
        if isinstance(child.tag, str)
        and etree.QName(child).localname in CATEGORY_SOURCE_NAMES
    ]
    return (category, sources[0]) if len(sources) == 1 else (category, None)


def cached_category_labels(series, point_budget=None):
    category, source = category_source(series)
    if category is None:
        return []
    if source is None:
        return None
    source_name = etree.QName(source).localname

    if source_name == "multiLvlStrRef":
        caches = source.findall(qn("c:multiLvlStrCache"))
        if len(caches) != 1 or (count := cache_point_count(caches[0])) is None:
            return None
        level_nodes = caches[0].findall(qn("c:lvl"))
        if len(level_nodes) > MAX_CATEGORY_LEVELS or (count and not level_nodes):
            return None
        levels = [cached_text_point_map(level, count) for level in level_nodes]
        if any(points is None for points in levels):
            return None
        if not consume_point_budget(point_budget, count * max(len(level_nodes), 1)):
            return None
        expanded_levels = []
        for level_index, points in enumerate(levels):
            if level_index == 0:
                expanded_levels.append([points.get(index, "") for index in range(count)])
                continue
            ordered = sorted(points.items())
            if not ordered:
                expanded_levels.append([""] * count)
                continue
            cursor = 0
            expanded = []
            for leaf_index in range(count):
                while cursor + 1 < len(ordered) and ordered[cursor + 1][0] <= leaf_index:
                    cursor += 1
                expanded.append(ordered[cursor][1])
            expanded_levels.append(expanded)
        return [
            [level[index] for level in reversed(expanded_levels)]
            for index in range(count)
        ]

    cache_name = {"strRef": "c:strCache", "numRef": "c:numCache"}.get(source_name)
    if cache_name is None:
        cache = source
    else:
        caches = source.findall(qn(cache_name))
        if len(caches) != 1:
            return None
        cache = caches[0]
    if (count := cache_point_count(cache)) is None:
        return None
    values = cached_text_point_map(cache, count)
    if values is None:
        return None
    if not consume_point_budget(point_budget, count):
        return None
    return [[values.get(index, "")] for index in range(count)]


def category_content(series, point_budget=None):
    labels = cached_category_labels(series, point_budget)
    _, source = category_source(series)
    formula = None if source is None else source.find(qn("c:f"))
    return {
        "categories": labels,
        "category_formula": formula.text if formula is not None else None,
        "category_cache_status": "unavailable" if labels is None else "available",
    }


def series_name_content(series):
    titles = series._element.xpath("./c:tx")
    if not titles:
        return {"name": ""}
    literal = titles[0].find(qn("c:v"))
    if literal is not None:
        return {"name": literal.text or ""}
    reference = titles[0].find(qn("c:strRef"))
    if reference is None:
        return {"name": None, "name_cache_status": "unavailable"}
    cache = reference.find(qn("c:strCache"))
    if cache is None:
        return {"name": None, "name_cache_status": "unavailable"}
    point_count = cache.find(qn("c:ptCount"))
    points = cache.findall(qn("c:pt"))
    if (point_count is None or point_count.get("val") != "1"
            or len(points) != 1 or points[0].get("idx") != "0"):
        return {"name": None, "name_cache_status": "unavailable"}
    value = points[0].find(qn("c:v"))
    if value is None:
        return {"name": None, "name_cache_status": "unavailable"}
    return {"name": value.text or ""}


def series_content(series, *, include_categories=False, point_budget=None):
    if point_budget is None:
        point_budget = {"remaining": MAX_CHART_POINTS}
    name_content = series_name_content(series)
    content = {
        **name_content,
        **(category_content(series, point_budget) if include_categories else {}),
    }
    x_source = getattr(series._element, "xVal", None)
    if x_source is None:
        value_source = getattr(series._element, "val", None)
        value_cache = cached_numeric_points(
            value_source, fill_missing=True, include_count=True,
            point_budget=point_budget,
        )
        if value_cache is None:
            return {**content, "values": None, "cache_status": "unavailable"}
        return {**content, "values": [value for _, value in value_cache["points"]]}
    x_cache = cached_numeric_points(
        x_source, include_count=True, point_budget=point_budget
    )
    y_cache = cached_numeric_points(
        getattr(series._element, "yVal", None), include_count=True,
        point_budget=point_budget,
    )
    if x_cache is None or y_cache is None:
        return {**content, "points": None, "cache_status": "unavailable"}
    content.update({
        "x_points": x_cache["points"], "x_point_count": x_cache["point_count"],
        "y_points": y_cache["points"], "y_point_count": y_cache["point_count"],
    })
    size_source = getattr(series._element, "bubbleSize", None)
    if size_source is not None:
        bubble_cache = cached_numeric_points(
            size_source, include_count=True, point_budget=point_budget
        )
        if bubble_cache is None:
            return {**content, "points": None, "cache_status": "unavailable"}
        content.update({
            "bubble_points": bubble_cache["points"],
            "bubble_point_count": bubble_cache["point_count"],
        })
    return content


def title_text(title):
    if title is None:
        return ""
    text_nodes = title.findall(qn("c:tx"))
    if len(text_nodes) != 1:
        return None
    sources = [child for child in text_nodes[0] if isinstance(child.tag, str)]
    if len(sources) != 1:
        return None
    if sources[0].tag == qn("c:rich"):
        paragraphs = sources[0].findall(qn("a:p"))
        return None if not paragraphs else "\n".join(
            paragraph.text for paragraph in paragraphs
        )
    if sources[0].tag != qn("c:strRef"):
        return None
    formulas = sources[0].findall(qn("c:f"))
    if len(formulas) != 1 or not formulas[0].text:
        return None
    caches = sources[0].findall(qn("c:strCache"))
    if (len(caches) != 1 or (count := cache_point_count(caches[0])) is None
            or count != 1):
        return None
    values = cached_text_point_map(caches[0], count)
    return None if values is None or set(values) != {0} else values[0]


def chart_axis_text(axis):
    return title_text(axis.find(qn("c:title")))


def chart_title_text(chart):
    titles = chart._element.xpath("./c:chart/c:title")
    return title_text(titles[0]) if len(titles) == 1 else ("" if not titles else None)


def chart_axes(chart):
    axes = chart._element.xpath(
        "./c:chart/c:plotArea/c:catAx | ./c:chart/c:plotArea/c:dateAx | "
        "./c:chart/c:plotArea/c:valAx | ./c:chart/c:plotArea/c:serAx"
    )
    return [
        {
            "kind": etree.QName(axis).localname,
            "id": axis.find(qn("c:axId")).get("val") if axis.find(qn("c:axId")) is not None else None,
            "position": (
                axis.find(qn("c:axPos")).get("val")
                if axis.find(qn("c:axPos")) is not None else None
            ),
            "cross_axis_id": (
                axis.find(qn("c:crossAx")).get("val")
                if axis.find(qn("c:crossAx")) is not None else None
            ),
            "title": chart_axis_text(axis),
        }
        for axis in axes
    ]


DIAGRAM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
SAFE_DIAGRAM_XML = etree.XMLParser(load_dtd=False, resolve_entities=False, no_network=True)


def smartart_content(shape):
    graphic_data = shape._element.find(".//" + qn("a:graphicData"))
    if graphic_data is None or graphic_data.get("uri") != DIAGRAM_NS:
        return None
    rel_ids = graphic_data.find(f".//{{{DIAGRAM_NS}}}relIds")
    relationship_id = None if rel_ids is None else rel_ids.get(qn("r:dm"))
    if not relationship_id:
        return {"name": shape.name, "status": "unreadable", "reason": "missing data relationship"}
    try:
        data_part = shape.part.related_part(relationship_id)
    except (KeyError, ValueError):
        return {"name": shape.name, "status": "unreadable", "reason": relationship_id}
    try:
        root = etree.fromstring(data_part.blob, parser=SAFE_DIAGRAM_XML)
    except etree.XMLSyntaxError as error:
        return {"name": shape.name, "status": "unreadable", "reason": str(error)}
    labels = [node.text for node in root.iter(qn("a:t")) if node.text]
    return {"name": shape.name, "status": "ok", "text": labels}


def extract_slide_content(slide, point_budget=None):
    shapes = list(iter_shapes(slide.shapes))
    if point_budget is None:
        point_budget = {"remaining": MAX_CHART_POINTS}
    text = slide_text_content(slide)
    tables = [
        table_cells(sh.table)
        for sh in shapes if sh.has_table
    ]
    charts = []
    for sh in shapes:
        if not sh.has_chart:
            continue
        chart = sh.chart
        chart_title = chart_title_text(chart)
        plots = []
        for plot in chart.plots:
            items = list(plot.series)
            has_xy_values = any(getattr(item._element, "xVal", None) is not None for item in items)
            series = [
                series_content(
                    item, include_categories=not has_xy_values, point_budget=point_budget
                )
                for item in items
            ]
            plots.append({
                "kind": type(plot).__name__,
                "series": series,
            })
        charts.append({
            "title": chart_title,
            "axes": chart_axes(chart),
            "plots": plots,
        })
    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
    pictures = [
        info for shape in shapes
        if (info := picture_content(shape)) is not None
    ]
    smartart = [
        info for shape in shapes
        if (info := smartart_content(shape)) is not None
    ]
    return {"text": text, "tables": tables, "charts": charts,
            "pictures": pictures, "smartart": smartart, "notes": notes}


content = extract_slide_content(open_validated_presentation("input.pptx").slides[0])
check("content inventory emits body text",
      any("old wording" in value["text"] for value in content["text"]), content)
check("content inventory emits table cell text",
      content["tables"][0][0][1]["text"] == "old cell text", content["tables"])
check(
    "content inventory emits chart title, categories, series, and values",
    content["charts"][0]["title"] == "Units by region"
    and content["charts"][0]["plots"][0]["series"][0]["categories"]
    == [["EU"], ["US"]]
    and content["charts"][0]["plots"][0]["series"][0]["category_cache_status"]
    == "available"
    and content["charts"][0]["plots"][0]["series"][0]["values"] == [120.0, 80.0],
    content["charts"],
)
check("content inventory emits raw category and value axis metadata",
      {axis["title"] for axis in content["charts"][0]["axes"]}
      == {"Region", "Units sold"}
      and all(axis["id"] and axis["position"] and axis["cross_axis_id"]
              for axis in content["charts"][0]["axes"]), content["charts"][0])

chart_title_prs = Presentation("input.pptx")
chart_title_chart = next(
    shape.chart for shape in chart_title_prs.slides[0].shapes if shape.has_chart
)
chart_title_element = chart_title_chart._element.xpath("./c:chart/c:title")[0]
chart_title_tx = chart_title_element.find(qn("c:tx"))
for child in list(chart_title_tx):
    chart_title_tx.remove(child)
chart_title_ref = OxmlElement("c:strRef")
chart_title_formula = OxmlElement("c:f")
chart_title_formula.text = "Sheet1!$G$1"
chart_title_cache = OxmlElement("c:strCache")
chart_title_count = OxmlElement("c:ptCount")
chart_title_count.set("val", "1")
chart_title_point = OxmlElement("c:pt")
chart_title_point.set("idx", "0")
chart_title_value = OxmlElement("c:v")
chart_title_value.text = "Cached chart title"
chart_title_point.append(chart_title_value)
chart_title_cache.extend([chart_title_count, chart_title_point])
chart_title_ref.extend([chart_title_formula, chart_title_cache])
chart_title_tx.append(chart_title_ref)
chart_title_prs.save("chart-title-strref.pptx")

chart_title_reopened = Presentation("chart-title-strref.pptx")
raw_title_chart = next(
    shape.chart for shape in chart_title_reopened.slides[0].shapes if shape.has_chart
)
raw_title_element = raw_title_chart._element.xpath("./c:chart/c:title")[0]
title_xml_before = etree.tostring(raw_title_element)
cached_chart_title = chart_title_text(raw_title_chart)
title_xml_after = etree.tostring(raw_title_element)
check("raw chart title reads a worksheet cache without mutating strRef",
      cached_chart_title == "Cached chart title"
      and title_xml_before == title_xml_after
      and raw_title_element.xpath("./c:tx/c:strRef/c:f")[0].text == "Sheet1!$G$1"
      and not raw_title_element.xpath("./c:tx/c:rich"),
      cached_chart_title)
inventory_title_xml_before = etree.tostring(raw_title_element)
cached_title_inventory = extract_slide_content(chart_title_reopened.slides[0])
inventory_title_xml_after = etree.tostring(raw_title_element)
check("full content inventory keeps a worksheet-backed chart title and its XML intact",
      cached_title_inventory["charts"][0]["title"] == "Cached chart title"
      and inventory_title_xml_before == inventory_title_xml_after
      and raw_title_element.xpath("./c:tx/c:strRef/c:f")[0].text == "Sheet1!$G$1"
      and not raw_title_element.xpath("./c:tx/c:rich"),
      cached_title_inventory["charts"][0])

cacheless_chart_title = copy.deepcopy(raw_title_element)
cacheless_chart_title.xpath("./c:tx/c:strRef")[0].remove(
    cacheless_chart_title.xpath("./c:tx/c:strRef/c:strCache")[0]
)
empty_chart_title = copy.deepcopy(raw_title_element)
empty_chart_title.xpath("./c:tx/c:strRef/c:strCache/c:pt/c:v")[0].text = None
ambiguous_chart_title = copy.deepcopy(raw_title_element)
ambiguous_chart_title.find(qn("c:tx")).append(OxmlElement("c:v"))
literal_only_chart_title = copy.deepcopy(raw_title_element)
literal_only_tx = literal_only_chart_title.find(qn("c:tx"))
literal_only_tx.remove(literal_only_tx.find(qn("c:strRef")))
literal_only_value = OxmlElement("c:v")
literal_only_value.text = "invalid literal title"
literal_only_tx.append(literal_only_value)
foreign_namespace_title = copy.deepcopy(raw_title_element)
foreign_namespace_tx = foreign_namespace_title.find(qn("c:tx"))
foreign_namespace_ref = foreign_namespace_tx.find(qn("c:strRef"))
foreign_namespace_ref.tag = "{urn:foreign-chart-title}strRef"
empty_rich_title = copy.deepcopy(raw_title_element)
empty_rich_tx = empty_rich_title.find(qn("c:tx"))
empty_rich_tx.remove(empty_rich_tx.find(qn("c:strRef")))
empty_rich_tx.append(OxmlElement("c:rich"))
duplicate_choice_title = copy.deepcopy(raw_title_element)
duplicate_choice_tx = duplicate_choice_title.find(qn("c:tx"))
duplicate_choice_tx.append(copy.deepcopy(duplicate_choice_tx.find(qn("c:strRef"))))
duplicate_cache_title = copy.deepcopy(raw_title_element)
duplicate_cache_ref = duplicate_cache_title.xpath("./c:tx/c:strRef")[0]
duplicate_cache_ref.append(copy.deepcopy(duplicate_cache_ref.find(qn("c:strCache"))))
missing_formula_title = copy.deepcopy(raw_title_element)
missing_formula_ref = missing_formula_title.xpath("./c:tx/c:strRef")[0]
missing_formula_ref.remove(missing_formula_ref.find(qn("c:f")))
duplicate_formula_title = copy.deepcopy(raw_title_element)
duplicate_formula_ref = duplicate_formula_title.xpath("./c:tx/c:strRef")[0]
duplicate_formula_ref.append(copy.deepcopy(duplicate_formula_ref.find(qn("c:f"))))
missing_count_title = copy.deepcopy(raw_title_element)
missing_count_cache = missing_count_title.xpath("./c:tx/c:strRef/c:strCache")[0]
missing_count_cache.remove(missing_count_cache.find(qn("c:ptCount")))
missing_point_title = copy.deepcopy(raw_title_element)
missing_point_cache = missing_point_title.xpath("./c:tx/c:strRef/c:strCache")[0]
missing_point_cache.remove(missing_point_cache.find(qn("c:pt")))
missing_value_title = copy.deepcopy(raw_title_element)
missing_value_point = missing_value_title.xpath("./c:tx/c:strRef/c:strCache/c:pt")[0]
missing_value_point.remove(missing_value_point.find(qn("c:v")))
duplicate_value_title = copy.deepcopy(raw_title_element)
duplicate_value_point = duplicate_value_title.xpath("./c:tx/c:strRef/c:strCache/c:pt")[0]
duplicate_value = OxmlElement("c:v")
duplicate_value.text = "second value"
duplicate_value_point.append(duplicate_value)

rich_paragraph_title = copy.deepcopy(raw_title_element)
rich_paragraph_tx = rich_paragraph_title.find(qn("c:tx"))
rich_paragraph_tx.remove(rich_paragraph_tx.find(qn("c:strRef")))
rich_title = OxmlElement("c:rich")
for text_value in ("First paragraph", "Second paragraph"):
    paragraph = OxmlElement("a:p")
    run = OxmlElement("a:r")
    text_node = OxmlElement("a:t")
    text_node.text = text_value
    run.append(text_node)
    paragraph.append(run)
    rich_title.append(paragraph)
rich_paragraph_tx.append(rich_title)
check("chart title cache parser distinguishes unavailable, empty, and ambiguous choices",
      title_text(cacheless_chart_title) is None
      and title_text(empty_chart_title) == ""
      and title_text(ambiguous_chart_title) is None
      and title_text(literal_only_chart_title) is None
      and title_text(foreign_namespace_title) is None
      and title_text(empty_rich_title) is None
      and title_text(duplicate_choice_title) is None
      and title_text(duplicate_cache_title) is None
      and title_text(missing_formula_title) is None
      and title_text(duplicate_formula_title) is None
      and title_text(missing_count_title) is None
      and title_text(missing_point_title) is None
      and title_text(missing_value_title) is None
      and title_text(duplicate_value_title) is None
      and title_text(rich_paragraph_title) == "First paragraph\nSecond paragraph")


def install_inherited_text(target_shapes, scratch_slide, text, *, grouped=False, hidden=False):
    if grouped:
        source_shape = scratch_slide.shapes.add_group_shape()
        group_texts = (text,) if isinstance(text, str) else tuple(text)
        for index, text_value in enumerate(group_texts):
            child = source_shape.shapes.add_textbox(
                Inches(1), Inches(1 + index), Inches(5), Inches(0.5)
            )
            child.text = text_value
    else:
        source_shape = scratch_slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(0.5)
        )
        source_shape.text = text
    copied = copy.deepcopy(source_shape._element)
    next_shape_id = target_shapes._next_shape_id
    for non_visual_properties in copied.iter(qn("p:cNvPr")):
        non_visual_properties.set("id", str(next_shape_id))
        next_shape_id += 1
    if hidden:
        copied.find(".//" + qn("p:cNvPr")).set("hidden", "1")
    target_shapes._spTree.insert_element_before(copied, "p:extLst")
    scratch_slide.shapes._spTree.remove(source_shape._element)


inheritance_prs = Presentation()
inheritance_layout = inheritance_prs.slide_layouts[5]
inheritance_slide = inheritance_prs.slides.add_slide(inheritance_layout)
inheritance_slide.shapes.title.text = "Actual slide title"
inheritance_layout_title = next(
    shape for shape in inheritance_layout.placeholders
    if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
)
inheritance_layout_title.text = "TEMPLATE PLACEHOLDER PROMPT"
inheritance_master_placeholder = next(iter(inheritance_layout.slide_master.placeholders))
inheritance_master_placeholder.text = "MASTER TEMPLATE PLACEHOLDER PROMPT"
install_inherited_text(
    inheritance_layout.slide_master.shapes, inheritance_slide, "Master disclaimer"
)
install_inherited_text(inheritance_layout.shapes, inheritance_slide, "Layout disclaimer")
install_inherited_text(
    inheritance_layout.shapes, inheritance_slide,
    ("Grouped layout first", "Grouped layout second"), grouped=True,
)
install_inherited_text(
    inheritance_layout.slide_master.shapes, inheritance_slide,
    "Hidden master copy", hidden=True,
)
install_inherited_text(
    inheritance_layout.shapes, inheritance_slide,
    "Hidden grouped layout copy", grouped=True, hidden=True,
)
install_inherited_text(
    inheritance_layout.shapes, inheritance_slide, "Whitespace-hidden layout copy"
)
whitespace_hidden_shape = next(
    shape for shape in inheritance_layout.shapes
    if shape.has_text_frame and shape.text_frame.text == "Whitespace-hidden layout copy"
)
whitespace_hidden_shape._element.find(".//" + qn("p:cNvPr")).set(
    "hidden", " \ttrue\r\n"
)
hidden_slide_shape = inheritance_slide.shapes.add_textbox(
    Inches(1), Inches(6), Inches(5), Inches(0.5)
)
hidden_slide_shape.text = "Hidden slide copy"
hidden_slide_shape._element.find(".//" + qn("p:cNvPr")).set("hidden", "1")
inheritance_prs.save("inherited-text.pptx")

inherited_slide = open_validated_presentation("inherited-text.pptx").slides[0]
inherited_text = extract_slide_content(inherited_slide)["text"]
inherited_pairs = [(item["source"], item["text"]) for item in inherited_text]
inherited_master_ids = [
    node.get("id") for node in inherited_slide.slide_layout.slide_master._element.xpath(
        ".//p:cNvPr"
    )
]
inherited_layout_ids = [
    node.get("id") for node in inherited_slide.slide_layout._element.xpath(".//p:cNvPr")
]
check("inherited-text fixture keeps unique non-visual shape ids per part",
      len(inherited_master_ids) == len(set(inherited_master_ids))
      and len(inherited_layout_ids) == len(set(inherited_layout_ids)),
      (inherited_master_ids, inherited_layout_ids))
check("content inventory includes visible master and layout text with provenance",
      inherited_pairs == [
          ("master", "Master disclaimer"),
          ("layout", "Layout disclaimer"),
          ("layout", "Grouped layout first"),
          ("layout", "Grouped layout second"),
          ("slide", "Actual slide title"),
      ],
      inherited_pairs)

slide_hidden_prs = Presentation("inherited-text.pptx")
slide_hidden_prs.slides[0]._element.set("showMasterSp", "0")
slide_hidden_prs.save("slide-hides-inherited-text.pptx")
slide_hidden_pairs = [
    (item["source"], item["text"])
    for item in extract_slide_content(
        open_validated_presentation("slide-hides-inherited-text.pptx").slides[0]
    )["text"]
]
check("slide showMasterSp=false hides both layout and master copy",
      ("slide", "Actual slide title") in slide_hidden_pairs
      and not any(source in {"layout", "master"} for source, _ in slide_hidden_pairs),
      slide_hidden_pairs)

layout_hidden_prs = Presentation("inherited-text.pptx")
layout_hidden_prs.slides[0].slide_layout._element.set("showMasterSp", "false")
layout_hidden_prs.save("layout-hides-master-text.pptx")
layout_hidden_pairs = [
    (item["source"], item["text"])
    for item in extract_slide_content(
        open_validated_presentation("layout-hides-master-text.pptx").slides[0]
    )["text"]
]
check("layout showMasterSp=false hides only master copy",
      ("layout", "Layout disclaimer") in layout_hidden_pairs
      and ("slide", "Actual slide title") in layout_hidden_pairs
      and not any(source == "master" for source, _ in layout_hidden_pairs),
      layout_hidden_pairs)

whitespace_false_prs = Presentation("inherited-text.pptx")
whitespace_false_prs.slides[0]._element.set("showMasterSp", " \t0\r\n")
whitespace_false_prs.save("whitespace-false-inherited-text.pptx")
whitespace_false_text = extract_slide_content(
    open_validated_presentation("whitespace-false-inherited-text.pptx").slides[0]
)["text"]
check("xsd whitespace around false hides inherited layers",
      all(item["source"] == "slide" for item in whitespace_false_text),
      whitespace_false_text)

whitespace_true_prs = Presentation("inherited-text.pptx")
whitespace_true_prs.slides[0]._element.set("showMasterSp", "\ttrue\n")
whitespace_true_prs.save("whitespace-true-inherited-text.pptx")
whitespace_true_pairs = [
    (item["source"], item["text"])
    for item in extract_slide_content(
        open_validated_presentation("whitespace-true-inherited-text.pptx").slides[0]
    )["text"]
]
check("xsd whitespace around true preserves inherited layers",
      ("master", "Master disclaimer") in whitespace_true_pairs
      and ("layout", "Layout disclaimer") in whitespace_true_pairs,
      whitespace_true_pairs)

invalid_visibility_values = (
    "maybe", "on", "off", "yes", "no", "TRUE", "False", "true false",
)
invalid_visibility_rejections = []
for invalid_value in invalid_visibility_values:
    invalid_visibility_slide = Presentation("inherited-text.pptx").slides[0]
    invalid_visibility_slide._element.set("showMasterSp", invalid_value)
    try:
        extract_slide_content(invalid_visibility_slide)
    except ValueError:
        invalid_visibility_rejections.append(invalid_value)
check("non-xsd showMasterSp values fail closed under optimized Python",
      invalid_visibility_rejections == list(invalid_visibility_values),
      invalid_visibility_rejections)

secondary_axis_prs = Presentation("input.pptx")
secondary_axis_chart = next(
    shape.chart for shape in secondary_axis_prs.slides[0].shapes
    if shape.has_chart and shape.chart.has_title
    and shape.chart.chart_title.text_frame.text == "Units by region"
)
plot_area = secondary_axis_chart._element.xpath("./c:chart/c:plotArea")[0]
primary_category_axis = plot_area.find(qn("c:catAx"))
primary_value_axis = plot_area.find(qn("c:valAx"))
primary_category_id = primary_category_axis.find(qn("c:axId")).get("val")
primary_value_id = primary_value_axis.find(qn("c:axId")).get("val")
secondary_category_id, secondary_value_id = "91000001", "91000002"

secondary_plot = copy.deepcopy(plot_area.xpath("./c:barChart")[0])
for axis_id in secondary_plot.findall(qn("c:axId")):
    axis_id.set("val", (
        secondary_category_id if axis_id.get("val") == primary_category_id
        else secondary_value_id
    ))
first_axis_index = min(plot_area.index(primary_category_axis), plot_area.index(primary_value_axis))
plot_area.insert(first_axis_index, secondary_plot)

secondary_category_axis = copy.deepcopy(primary_category_axis)
secondary_category_axis.find(qn("c:axId")).set("val", secondary_category_id)
secondary_category_axis.find(qn("c:axPos")).set("val", "t")
secondary_category_axis.find(qn("c:crossAx")).set("val", secondary_value_id)
secondary_category_text = secondary_category_axis.xpath("./c:title//a:t")[0]
secondary_category_text.text = "Secondary"
secondary_category_run = secondary_category_text.getparent()
secondary_category_break = OxmlElement("a:br")
secondary_category_tail = OxmlElement("a:r")
secondary_category_tail_text = OxmlElement("a:t")
secondary_category_tail_text.text = "region"
secondary_category_tail.append(secondary_category_tail_text)
secondary_category_run.addnext(secondary_category_break)
secondary_category_break.addnext(secondary_category_tail)

secondary_value_axis = copy.deepcopy(primary_value_axis)
secondary_value_axis.find(qn("c:axId")).set("val", secondary_value_id)
secondary_value_axis.find(qn("c:axPos")).set("val", "r")
secondary_value_axis.find(qn("c:crossAx")).set("val", secondary_category_id)
secondary_title_tx = secondary_value_axis.find(qn("c:title")).find(qn("c:tx"))
for child in list(secondary_title_tx):
    secondary_title_tx.remove(child)
secondary_title_ref = OxmlElement("c:strRef")
secondary_title_formula = OxmlElement("c:f")
secondary_title_formula.text = "Sheet1!$F$1"
secondary_title_cache = OxmlElement("c:strCache")
secondary_title_count = OxmlElement("c:ptCount")
secondary_title_count.set("val", "1")
secondary_title_point = OxmlElement("c:pt")
secondary_title_point.set("idx", "0")
secondary_title_value = OxmlElement("c:v")
secondary_title_value.text = "Percent"
secondary_title_point.append(secondary_title_value)
secondary_title_cache.extend([secondary_title_count, secondary_title_point])
secondary_title_ref.extend([secondary_title_formula, secondary_title_cache])
secondary_title_tx.append(secondary_title_ref)
empty_secondary_axis = copy.deepcopy(secondary_value_axis)
empty_secondary_axis.xpath("./c:title/c:tx/c:strRef/c:strCache/c:pt/c:v")[0].text = None
check("raw worksheet-backed axis title preserves an explicit empty cache value",
      chart_axis_text(empty_secondary_axis) == "", chart_axis_text(empty_secondary_axis))
multi_point_axis = copy.deepcopy(secondary_value_axis)
multi_point_axis.xpath("./c:title/c:tx/c:strRef/c:strCache/c:ptCount")[0].set("val", "2")
check("worksheet-backed axis title rejects a non-scalar cache before expansion",
      chart_axis_text(multi_point_axis) is None)

plot_area.extend([secondary_category_axis, secondary_value_axis])
secondary_axis_prs.save("secondary-axes.pptx")
secondary_axis_reopened = Presentation("secondary-axes.pptx")
secondary_axis_content = extract_slide_content(
    secondary_axis_reopened.slides[0]
)["charts"][0]
secondary_axes = {axis["title"]: axis for axis in secondary_axis_content["axes"]}
secondary_axis_reopened_chart = next(
    shape.chart for shape in secondary_axis_reopened.slides[0].shapes
    if shape.has_chart and shape.chart.has_title
    and shape.chart.chart_title.text_frame.text == "Units by region"
)
secondary_title_refs = secondary_axis_reopened_chart._element.xpath(
    f'./c:chart/c:plotArea/c:valAx[c:axId[@val="{secondary_value_id}"]]'
    '/c:title/c:tx/c:strRef'
)
check("combination-chart inventory includes all primary and secondary axes",
      len(secondary_axis_content["plots"]) == 2
      and set(secondary_axes) == {"Region", "Units sold", "Secondary\vregion", "Percent"}
      and secondary_axes["Secondary\vregion"] == {
          "kind": "catAx", "id": secondary_category_id, "position": "t",
          "cross_axis_id": secondary_value_id, "title": "Secondary\vregion",
      }
      and secondary_axes["Percent"] == {
          "kind": "valAx", "id": secondary_value_id, "position": "r",
          "cross_axis_id": secondary_category_id, "title": "Percent",
      }, secondary_axis_content)
check("raw axis-title inventory preserves a worksheet-backed strRef",
      len(secondary_title_refs) == 1
      and secondary_title_refs[0].find(qn("c:f")).text == "Sheet1!$F$1",
      secondary_title_refs)
chart_by_title = {item["title"]: item for item in content["charts"]}
check("content inventory emits XY x/y points",
      chart_by_title["XY trend"]["plots"][0]["series"][0]["x_points"]
      == [(0, 1.0), (1, 3.0)]
      and chart_by_title["XY trend"]["plots"][0]["series"][0]["x_point_count"] == 2
      and chart_by_title["XY trend"]["plots"][0]["series"][0]["y_points"]
      == [(0, 2.0), (1, 4.0)]
      and chart_by_title["XY trend"]["plots"][0]["series"][0]["y_point_count"] == 2)
check("content inventory emits bubble x/y/size points",
      chart_by_title["Bubble risk"]["plots"][0]["series"][0]["x_points"] == [(0, 3.0)]
      and chart_by_title["Bubble risk"]["plots"][0]["series"][0]["y_points"] == [(0, 4.0)]
      and chart_by_title["Bubble risk"]["plots"][0]["series"][0]["bubble_points"]
      == [(0, 5.0)]
      and chart_by_title["Bubble risk"]["plots"][0]["series"][0]["bubble_point_count"]
      == 1)
check("content inventory emits notes text", "regional split" in content["notes"], content["notes"])

# Category/value series can retain an external workbook formula without a numCache.
category_prs = Presentation()
category_slide = category_prs.slides.add_slide(category_prs.slide_layouts[6])
category_data = ChartData()
category_data.categories = ["A", "B"]
category_data.add_series("Missing cache", (1, 2))
category_data.add_series("Cached series", (3, 4))
category_slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(0.5), Inches(6), Inches(4), category_data,
)
category_prs.save("category-cache-source.pptx")

shared_deck_budget = {"remaining": 10}
budget_source = Presentation("category-cache-source.pptx").slides[0]
first_budgeted_slide = extract_slide_content(budget_source, shared_deck_budget)
second_budgeted_slide = extract_slide_content(budget_source, shared_deck_budget)
first_budgeted_items = first_budgeted_slide["charts"][0]["plots"][0]["series"]
second_budgeted_items = second_budgeted_slide["charts"][0]["plots"][0]["series"]
check("one shared logical-point budget bounds aggregate inventory across slides",
      all(item["values"] is not None for item in first_budgeted_items)
      and second_budgeted_items[0]["categories"] == [["A"], ["B"]]
      and second_budgeted_items[0]["values"] is None
      and second_budgeted_items[1]["categories"] is None
      and shared_deck_budget["remaining"] == 0,
      (first_budgeted_items, second_budgeted_items, shared_deck_budget))

series_name_mutated = Presentation("category-cache-source.pptx")
series_name_chart = next(
    shape.chart for shape in series_name_mutated.slides[0].shapes if shape.has_chart
)
series_name_ref = list(series_name_chart.plots[0].series)[0]._element.xpath(
    "./c:tx/c:strRef"
)[0]
series_name_formula = series_name_ref.find(qn("c:f")).text
series_name_ref.remove(series_name_ref.find(qn("c:strCache")))
series_name_mutated.save("series-name-cacheless.pptx")
series_name_reopened = Presentation("series-name-cacheless.pptx")
series_name_chart = next(
    shape.chart for shape in series_name_reopened.slides[0].shapes if shape.has_chart
)
series_name_series = list(series_name_chart.plots[0].series)[0]
series_name_result = series_content(series_name_series)
check("cacheless series title retains its worksheet formula",
      series_name_series._element.xpath("./c:tx/c:strRef/c:f")[0].text
      == series_name_formula)
check("cacheless series title is unavailable without aborting value inventory",
      series_name_result == {
          "name": None, "name_cache_status": "unavailable", "values": [1.0, 2.0],
      }, series_name_result)
series_name_inventory = extract_slide_content(series_name_reopened.slides[0])
series_name_inventory_items = series_name_inventory["charts"][0]["plots"][0]["series"]
check("deck inventory continues through cached siblings after an unavailable series title",
      series_name_inventory_items[0]["name"] is None
      and series_name_inventory_items[0]["values"] == [1.0, 2.0]
      and series_name_inventory_items[1]["name"] == "Cached series"
      and series_name_inventory_items[1]["values"] == [3.0, 4.0],
      series_name_inventory_items)


class RaisingNameSeries:
    def __init__(self, element, values):
        self._element = element
        self.values = values

    @property
    def name(self):
        raise RuntimeError("series.name must not be accessed")


raising_name_result = series_content(RaisingNameSeries(
    series_name_series._element, list(series_name_series.values),
))
check("cacheless name path never touches the python-pptx name property",
      raising_name_result == series_name_result, raising_name_result)

source_name_series = list(series_name_chart.plots[0].series)[1]
malformed_name_element = copy.deepcopy(source_name_series._element)
malformed_cache = malformed_name_element.find(qn("c:tx")).find(
    qn("c:strRef")
).find(qn("c:strCache"))
for point in malformed_cache.findall(qn("c:pt")):
    malformed_cache.remove(point)
malformed_name_result = series_content(RaisingNameSeries(
    malformed_name_element, list(source_name_series.values),
))
check("incomplete series-name cache is unavailable while values remain readable",
      malformed_name_result == {
          "name": None, "name_cache_status": "unavailable", "values": [3.0, 4.0],
      }, malformed_name_result)

literal_name_element = copy.deepcopy(source_name_series._element)
literal_tx = literal_name_element.find(qn("c:tx"))
literal_tx.remove(literal_tx.find(qn("c:strRef")))
literal_value = OxmlElement("c:v")
literal_value.text = "Literal series"
literal_tx.append(literal_value)
literal_name_result = series_content(RaisingNameSeries(
    literal_name_element, list(source_name_series.values),
))
check("literal series title is read directly without the name property",
      literal_name_result["name"] == "Literal series", literal_name_result)

untitled_element = copy.deepcopy(source_name_series._element)
untitled_element.remove(untitled_element.find(qn("c:tx")))
untitled_result = series_content(RaisingNameSeries(
    untitled_element, list(source_name_series.values),
))
check("series without a title remains a readable unnamed series",
      untitled_result["name"] == "", untitled_result)

category_mutated = Presentation("category-cache-source.pptx")
category_chart = next(
    shape.chart for shape in category_mutated.slides[0].shapes if shape.has_chart
)
category_series = list(category_chart.plots[0].series)
missing_value_ref = category_series[0]._element.val.find(qn("c:numRef"))
missing_formula_before = missing_value_ref.find(qn("c:f")).text
missing_value_ref.remove(missing_value_ref.find(qn("c:numCache")))
category_mutated.save("category-cacheless.pptx")

category_reopened = Presentation("category-cacheless.pptx")
category_reopened_chart = next(
    shape.chart for shape in category_reopened.slides[0].shapes if shape.has_chart
)
category_reopened_series = list(category_reopened_chart.plots[0].series)
reopened_value_ref = category_reopened_series[0]._element.val.find(qn("c:numRef"))
reopened_formula = reopened_value_ref.find(qn("c:f"))
with zipfile.ZipFile("category-cacheless.pptx") as category_archive:
    embedded_workbook_remains = any(
        name.startswith("ppt/embeddings/") for name in category_archive.namelist()
    )
check(
    "cacheless category series retains its formula and embedded workbook",
    reopened_formula is not None
    and reopened_formula.text == missing_formula_before
    and reopened_value_ref.find(qn("c:numCache")) is None
    and embedded_workbook_remains,
)
category_inventory = extract_slide_content(category_reopened.slides[0])
category_inventory_series = category_inventory["charts"][0]["plots"][0]["series"]
check(
    "cacheless category series is unavailable without aborting inventory",
    category_inventory_series[0]["name"] == "Missing cache"
    and category_inventory_series[0]["values"] is None
    and category_inventory_series[0]["cache_status"] == "unavailable",
    category_inventory_series,
)
check(
    "cached category series still reports its values after a cacheless sibling",
    category_inventory_series[1]["name"] == "Cached series"
    and category_inventory_series[1]["values"] == [3.0, 4.0],
    category_inventory_series,
)

marker_element = copy.deepcopy(category_series[1]._element)
marker_values = marker_element.xpath("./c:val/c:numRef/c:numCache/c:pt")
marker_values[0].find(qn("c:v")).text = "#N/A"
marker_values[1].find(qn("c:v")).text = None
marker_series = type("CachedMarkerSeries", (), {"_element": marker_element})()
marker_result = series_content(marker_series, include_categories=True)
check("category series preserves #N/A and blank cache markers without series.values",
      marker_result["values"] == ["#N/A", None], marker_result)
omitted_marker_element = copy.deepcopy(marker_element)
omitted_marker_cache = omitted_marker_element.xpath("./c:val/c:numRef/c:numCache")[0]
omitted_marker_cache.remove(omitted_marker_cache.findall(qn("c:pt"))[1])
omitted_marker_series = type(
    "SparseCategorySeries", (), {"_element": omitted_marker_element}
)()
omitted_marker_result = series_content(omitted_marker_series, include_categories=True)
check("category series densifies an omitted sparse point as blank",
      omitted_marker_result["values"] == ["#N/A", None], omitted_marker_result)

divergent_prs = Presentation("category-cache-source.pptx")
divergent_chart = next(
    shape.chart for shape in divergent_prs.slides[0].shapes if shape.has_chart
)
divergent_series = list(divergent_chart.plots[0].series)
second_category = divergent_series[1]._element.find(qn("c:cat"))
second_category.find(qn("c:strRef")).find(qn("c:f")).text = "Sheet1!$D$2:$D$3"
second_points = second_category.xpath("./c:strRef/c:strCache/c:pt")
second_points[0].find(qn("c:v")).text = "North"
second_points[1].find(qn("c:v")).text = "South"
divergent_inventory = extract_slide_content(divergent_prs.slides[0])
divergent_items = divergent_inventory["charts"][0]["plots"][0]["series"]
check("each category series retains its own formula and cached labels",
      divergent_items[0]["categories"] == [["A"], ["B"]]
      and divergent_items[1]["categories"] == [["North"], ["South"]]
      and divergent_items[0]["category_formula"] != divergent_items[1]["category_formula"],
      divergent_items)
commented_category_element = copy.deepcopy(divergent_series[0]._element)
commented_category_element.find(qn("c:cat")).insert(0, etree.Comment("source follows"))
commented_category_series = type(
    "CommentedCategorySeries", (), {"_element": commented_category_element}
)()
check("category source selection ignores preserved XML comments",
      category_content(commented_category_series)["categories"] == [["A"], ["B"]])
malformed_category_element = copy.deepcopy(divergent_series[1]._element)
malformed_category_cache = malformed_category_element.xpath(
    "./c:cat/c:strRef/c:strCache"
)[0]
malformed_category_cache.append(copy.deepcopy(
    malformed_category_cache.findall(qn("c:pt"))[0]
))
malformed_category_series = type(
    "MalformedCategorySeries", (), {"_element": malformed_category_element}
)()
malformed_category_result = category_content(malformed_category_series)
check("duplicate category cache indexes fail closed for only that series",
      malformed_category_result["categories"] is None
      and malformed_category_result["category_cache_status"] == "unavailable",
      malformed_category_result)

duplicate_count_element = copy.deepcopy(divergent_series[1]._element)
duplicate_count_cache = duplicate_count_element.xpath("./c:cat/c:strRef/c:strCache")[0]
duplicate_count_cache.insert(1, copy.deepcopy(duplicate_count_cache.find(qn("c:ptCount"))))
duplicate_count_series = type(
    "DuplicateCategoryCountSeries", (), {"_element": duplicate_count_element}
)()
duplicate_count_budget = {"remaining": 2}
duplicate_count_result = category_content(duplicate_count_series, duplicate_count_budget)
valid_after_malformed = category_content(divergent_series[0], duplicate_count_budget)
check("duplicate category ptCount fails closed without starving a valid sibling",
      duplicate_count_result["categories"] is None
      and valid_after_malformed["categories"] == [["A"], ["B"]]
      and duplicate_count_budget["remaining"] == 0,
      (duplicate_count_result, valid_after_malformed, duplicate_count_budget))

choice_conflict_element = copy.deepcopy(divergent_series[1]._element)
choice_conflict_category = choice_conflict_element.find(qn("c:cat"))
choice_conflict_category.append(parse_xml(
    '<c:multiLvlStrRef xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">'
    '<c:f>Sheet1!$E$2:$E$3</c:f><c:multiLvlStrCache><c:ptCount val="2"/>'
    '<c:lvl><c:pt idx="0"><c:v>X</c:v></c:pt>'
    '<c:pt idx="1"><c:v>Y</c:v></c:pt></c:lvl>'
    '</c:multiLvlStrCache></c:multiLvlStrRef>'
))
choice_conflict_series = type(
    "ChoiceConflictSeries", (), {"_element": choice_conflict_element}
)()
choice_conflict_result = category_content(choice_conflict_series)
check("conflicting category source choices fail closed without mixing formula and labels",
      choice_conflict_result["categories"] is None
      and choice_conflict_result["category_formula"] is None
      and choice_conflict_result["category_cache_status"] == "unavailable",
      choice_conflict_result)

huge_category_element = copy.deepcopy(divergent_series[1]._element)
huge_category_element.xpath("./c:cat/c:strRef/c:strCache/c:ptCount")[0].set(
    "val", "4294967295"
)
huge_category_series = type(
    "HugeCategoryCountSeries", (), {"_element": huge_category_element}
)()
check("oversized logical category counts fail before allocation",
      category_content(huge_category_series)["categories"] is None)

aggregate_budget = {"remaining": 3}
first_budgeted_categories = category_content(divergent_series[0], aggregate_budget)
second_budgeted_categories = category_content(divergent_series[1], aggregate_budget)
check("shared category budget bounds aggregate dense inventory",
      first_budgeted_categories["category_cache_status"] == "available"
      and second_budgeted_categories["category_cache_status"] == "unavailable"
      and aggregate_budget["remaining"] == 1,
      (first_budgeted_categories, second_budgeted_categories, aggregate_budget))

multi_level_element = copy.deepcopy(divergent_series[0]._element)
multi_level_category = multi_level_element.find(qn("c:cat"))
multi_level_category.remove(next(iter(multi_level_category)))
multi_level_category.append(parse_xml(
    '<c:multiLvlStrRef xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">'
    '<c:f>Sheet1!$A$2:$A$5</c:f><c:multiLvlStrCache><c:ptCount val="4"/>'
    '<c:lvl><c:pt idx="0"><c:v>A</c:v></c:pt><c:pt idx="1"><c:v>B</c:v></c:pt>'
    '<c:pt idx="2"><c:v>C</c:v></c:pt><c:pt idx="3"><c:v>D</c:v></c:pt></c:lvl>'
    '<c:lvl><c:pt idx="0"><c:v>G1</c:v></c:pt>'
    '<c:pt idx="2"><c:v>G2</c:v></c:pt></c:lvl>'
    '</c:multiLvlStrCache></c:multiLvlStrRef>'
))
multi_level_series = type("MultiLevelSeries", (), {"_element": multi_level_element})()
multi_level_result = category_content(multi_level_series)
check("multi-level categories flatten parent-to-child with linear carry-forward",
      multi_level_result["categories"]
      == [["G1", "A"], ["G1", "B"], ["G2", "C"], ["G2", "D"]]
      and multi_level_result["category_formula"] == "Sheet1!$A$2:$A$5",
      multi_level_result)

category_labels_mutated = Presentation("category-cache-source.pptx")
category_labels_chart = next(
    shape.chart for shape in category_labels_mutated.slides[0].shapes if shape.has_chart
)
category_ref = category_labels_chart.plots[0]._element.xpath(
    "./c:ser[1]/c:cat/c:strRef"
)[0]
category_formula_before = category_ref.find(qn("c:f")).text
category_ref.remove(category_ref.find(qn("c:strCache")))
category_labels_mutated.save("category-label-cacheless.pptx")
category_labels_reopened = Presentation("category-label-cacheless.pptx")
category_labels_chart = next(
    shape.chart for shape in category_labels_reopened.slides[0].shapes if shape.has_chart
)
category_ref = category_labels_chart.plots[0]._element.xpath(
    "./c:ser[1]/c:cat/c:strRef"
)[0]
with zipfile.ZipFile("category-label-cacheless.pptx") as category_archive:
    category_workbook_remains = any(
        name.startswith("ppt/embeddings/") for name in category_archive.namelist()
    )
category_labels_inventory = extract_slide_content(category_labels_reopened.slides[0])
category_plot_inventory = category_labels_inventory["charts"][0]["plots"][0]
check("cacheless categories retain their worksheet formula and embedded workbook",
      category_ref.find(qn("c:f")).text == category_formula_before
      and category_ref.find(qn("c:strCache")) is None
      and category_workbook_remains)
check("missing category cache is reported without aborting later chart inventory",
      category_plot_inventory["series"][0]["categories"] is None
      and category_plot_inventory["series"][0]["category_cache_status"] == "unavailable"
      and category_plot_inventory["series"][1]["categories"] == [["A"], ["B"]]
      and category_plot_inventory["series"][1]["values"] == [3.0, 4.0],
      category_plot_inventory)

diagram_frame = etree.fromstring(f'''<p:graphicFrame
    xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
    xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
    xmlns:dgm="{DIAGRAM_NS}"
    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <a:graphic><a:graphicData uri="{DIAGRAM_NS}"><dgm:relIds r:dm="rIdSmart"/></a:graphicData></a:graphic>
</p:graphicFrame>'''.encode())
diagram_data = f'''<dgm:dataModel xmlns:dgm="{DIAGRAM_NS}"
    xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <dgm:ptLst><dgm:pt><dgm:t><a:p><a:r><a:t>Revenue</a:t></a:r></a:p></dgm:t></dgm:pt></dgm:ptLst>
</dgm:dataModel>'''.encode()
diagram_part = type("DiagramPart", (), {"blob": diagram_data})()
diagram_owner = type("SlidePart", (), {
    "related_part": lambda self, relationship_id: {"rIdSmart": diagram_part}[relationship_id],
})()
diagram_shape = type("SmartArtShape", (), {
    "name": "SmartArt 1", "_element": diagram_frame, "part": diagram_owner,
})()
check("SmartArt inventory extracts text from the diagram data part",
      smartart_content(diagram_shape)
      == {"name": "SmartArt 1", "status": "ok", "text": ["Revenue"]})
missing_diagram_shape = type("SmartArtShape", (), {
    "name": "Broken SmartArt", "_element": diagram_frame,
    "part": type("SlidePart", (), {
        "related_part": lambda self, relationship_id: {}[relationship_id],
    })(),
})()
check("SmartArt inventory reports an unresolved diagram relationship",
      smartart_content(missing_diagram_shape)["status"] == "unreadable")

placeholder_png = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9Y9ZlRYAAAAASUVORK5CYII="
)
with open("placeholder-picture.png", "wb") as stream:
    stream.write(placeholder_png)
placeholder_prs = Presentation()
picture_layout = next(
    layout for layout in placeholder_prs.slide_layouts
    if any(
        shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
        for shape in layout.placeholders
    )
)
placeholder_slide = placeholder_prs.slides.add_slide(picture_layout)
picture_placeholder = next(
    shape for shape in placeholder_slide.placeholders
    if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
)
placeholder_picture = picture_placeholder.insert_picture("placeholder-picture.png")
placeholder_prs.save("picture-placeholder.pptx")
placeholder_reopened = Presentation("picture-placeholder.pptx").slides[0]
reopened_picture = next(
    shape for shape in placeholder_reopened.placeholders if hasattr(shape, "image")
)
placeholder_inventory = extract_slide_content(placeholder_reopened)["pictures"]
check("picture placeholder remains a placeholder after image insertion",
      reopened_picture.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER,
      reopened_picture.shape_type)
check("picture inventory includes populated picture placeholders and image metadata",
      len(placeholder_inventory) == 1
      and placeholder_inventory[0]["extension"] == "png"
      and placeholder_inventory[0]["bytes"] == len(placeholder_png),
      placeholder_inventory)

# A dangling r:embed must produce an unreadable record and allow later slides to inventory.
broken_picture_source = Presentation("picture-placeholder.pptx")
later_picture_slide = broken_picture_source.slides.add_slide(
    broken_picture_source.slide_layouts[6]
)
later_picture_slide.shapes.add_picture(
    "placeholder-picture.png", Inches(1), Inches(1), Inches(1), Inches(1)
)
broken_picture_source.save("broken-picture-source.pptx")
with zipfile.ZipFile("broken-picture-source.pptx") as source_archive:
    broken_picture_payload = {
        info.filename: source_archive.read(info) for info in source_archive.infolist()
    }
broken_slide_root = etree.fromstring(broken_picture_payload["ppt/slides/slide1.xml"])
broken_blips = broken_slide_root.findall(".//" + qn("a:blip"))
require(len(broken_blips) == 1, "broken-picture fixture expected one slide-1 image")
broken_relationship_id = broken_blips[0].get(qn("r:embed"))
relationships_path = "ppt/slides/_rels/slide1.xml.rels"
relationships_root = etree.fromstring(broken_picture_payload[relationships_path])
package_relationship_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
broken_relationships = [
    relation for relation in relationships_root.findall(
        f"{{{package_relationship_ns}}}Relationship"
    )
    if relation.get("Id") == broken_relationship_id
]
require(len(broken_relationships) == 1,
        "broken-picture fixture expected one matching relationship")
relationships_root.remove(broken_relationships[0])
broken_picture_payload[relationships_path] = etree.tostring(
    relationships_root, xml_declaration=True, encoding="UTF-8", standalone=True,
)
with zipfile.ZipFile("broken-picture-relationship.pptx", "w", zipfile.ZIP_DEFLATED) as archive:
    for member_name, member_data in broken_picture_payload.items():
        archive.writestr(member_name, member_data)
broken_picture_deck = open_validated_presentation("broken-picture-relationship.pptx")
broken_picture_inventory = [
    extract_slide_content(slide)["pictures"] for slide in broken_picture_deck.slides
]
check("broken picture relationship is explicit instead of aborting inventory",
      broken_picture_inventory[0] == [{
          "name": reopened_picture.name,
          "status": "unreadable",
          "relationship_id": broken_relationship_id,
          "reason": "missing or invalid image relationship or payload",
      }], broken_picture_inventory[0])
check("picture inventory continues to a healthy later slide after a broken relationship",
      len(broken_picture_inventory[1]) == 1
      and broken_picture_inventory[1][0]["extension"] == "png",
      broken_picture_inventory[1])

# Directly hidden leaf shapes and children of a hidden group must be absent from every
# non-text inventory, while visible siblings after them remain discoverable.
hidden_inventory_prs = Presentation()
hidden_inventory_slide = hidden_inventory_prs.slides.add_slide(
    hidden_inventory_prs.slide_layouts[6]
)

def set_shape_hidden(shape):
    properties = shape._element.find(".//" + qn("p:cNvPr"))
    require(properties is not None, f"shape {shape.name!r} has no cNvPr")
    properties.set("hidden", "1")

hidden_table_shape = hidden_inventory_slide.shapes.add_table(
    1, 1, Inches(0.2), Inches(0.2), Inches(1), Inches(0.5)
)
hidden_table_shape.table.cell(0, 0).text = "secret table"
set_shape_hidden(hidden_table_shape)

hidden_chart_data = ChartData()
hidden_chart_data.categories = ["secret"]
hidden_chart_data.add_series("secret", (9,))
hidden_chart_shape = hidden_inventory_slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1.3), Inches(0.2), Inches(1.5), Inches(1), hidden_chart_data,
)
set_shape_hidden(hidden_chart_shape)

hidden_picture_shape = hidden_inventory_slide.shapes.add_picture(
    "placeholder-picture.png", Inches(3), Inches(0.2), Inches(1), Inches(1)
)
set_shape_hidden(hidden_picture_shape)

def smartart_frame(shape_id, name, relationship_id, *, hidden=False):
    hidden_attribute = ' hidden="1"' if hidden else ""
    return etree.fromstring(f'''<p:graphicFrame
        xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
        xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
        xmlns:dgm="{DIAGRAM_NS}"
        xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:nvGraphicFramePr>
        <p:cNvPr id="{shape_id}" name="{name}"{hidden_attribute}/>
        <p:cNvGraphicFramePr/><p:nvPr/>
      </p:nvGraphicFramePr>
      <p:xfrm><a:off x="0" y="0"/><a:ext cx="914400" cy="914400"/></p:xfrm>
      <a:graphic><a:graphicData uri="{DIAGRAM_NS}">
        <dgm:relIds r:dm="{relationship_id}"/>
      </a:graphicData></a:graphic>
    </p:graphicFrame>'''.encode())

hidden_inventory_slide.shapes._spTree.append(
    smartart_frame(910, "Hidden SmartArt", "rIdHiddenSmart", hidden=True)
)

nested_hidden_table = hidden_inventory_slide.shapes.add_table(
    1, 1, Inches(4.2), Inches(0.2), Inches(1), Inches(0.5)
)
nested_hidden_table.table.cell(0, 0).text = "secret grouped table"
nested_hidden_element = nested_hidden_table._element
hidden_group_element = parse_xml(GRP)
hidden_group_element.find(".//" + qn("p:cNvPr")).set("hidden", "1")
nested_hidden_element.getparent().replace(nested_hidden_element, hidden_group_element)
hidden_group_element.append(nested_hidden_element)

visible_table_shape = hidden_inventory_slide.shapes.add_table(
    1, 1, Inches(0.2), Inches(2), Inches(1), Inches(0.5)
)
visible_table_shape.table.cell(0, 0).text = "visible table"
visible_chart_data = ChartData()
visible_chart_data.categories = ["visible"]
visible_chart_data.add_series("visible", (1,))
visible_chart_shape = hidden_inventory_slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1.3), Inches(2), Inches(1.5), Inches(1), visible_chart_data,
)
visible_picture_shape = hidden_inventory_slide.shapes.add_picture(
    "placeholder-picture.png", Inches(3), Inches(2), Inches(1), Inches(1)
)
hidden_inventory_slide.shapes._spTree.append(
    smartart_frame(911, "Visible SmartArt", "rIdVisibleSmart")
)
hidden_inventory_prs.save("hidden-nontext-inventory.pptx")
hidden_inventory_reopened = open_validated_presentation("hidden-nontext-inventory.pptx")
hidden_inventory_slide = hidden_inventory_reopened.slides[0]
hidden_raw_names = [shape.name for shape in hidden_inventory_slide.shapes]
hidden_visible_names = [shape.name for shape in iter_shapes(hidden_inventory_slide.shapes)]
hidden_inventory = extract_slide_content(hidden_inventory_slide)
check("hidden non-text fixture retains direct and group-hidden objects (negative control)",
      {"Hidden SmartArt", hidden_table_shape.name, hidden_chart_shape.name,
       hidden_picture_shape.name, "demo group"} <= set(hidden_raw_names),
      hidden_raw_names)
check("visibility walker prunes direct hidden leaves and every child of a hidden group",
      "Hidden SmartArt" not in hidden_visible_names
      and hidden_table_shape.name not in hidden_visible_names
      and hidden_chart_shape.name not in hidden_visible_names
      and hidden_picture_shape.name not in hidden_visible_names
      and "secret grouped table" not in str(hidden_inventory),
      (hidden_visible_names, hidden_inventory))
check("all non-text inventories retain only visible siblings after hidden shapes",
      len(hidden_inventory["tables"]) == 1
      and hidden_inventory["tables"][0][0][0]["text"] == "visible table"
      and len(hidden_inventory["charts"]) == 1
      and hidden_inventory["charts"][0]["plots"][0]["series"][0]["values"] == [1.0]
      and len(hidden_inventory["pictures"]) == 1
      and hidden_inventory["pictures"][0]["extension"] == "png"
      and hidden_inventory["smartart"] == [{
          "name": "Visible SmartArt", "status": "unreadable",
          "reason": "rIdVisibleSmart",
      }], hidden_inventory)

hidden_title_prs = Presentation()
hidden_title_slide = hidden_title_prs.slides.add_slide(hidden_title_prs.slide_layouts[5])
hidden_title_slide.shapes.title.text = "confidential hidden title"
set_shape_hidden(hidden_title_slide.shapes.title)
hidden_title_prs.save("hidden-title.pptx")
hidden_title_slide = open_validated_presentation("hidden-title.pptx").slides[0]
hidden_title_shape = hidden_title_slide.shapes.title
guarded_title = (
    hidden_title_shape.text_frame.text
    if hidden_title_shape is not None and not shape_is_hidden(hidden_title_shape) else ""
)
check("main slide summary does not expose a hidden title", guarded_title == "", guarded_title)

merged_prs = Presentation()
merged_slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])
merged_table = merged_slide.shapes.add_table(
    3, 3, Inches(1), Inches(1), Inches(6), Inches(3)
).table
merged_table.cell(0, 0).text = "Merged heading"
merged_table.cell(0, 0).merge(merged_table.cell(1, 1))
merged_table.cell(2, 0).text = "ordinary cell"
merged_prs.save("merged-table.pptx")
merged_inventory = extract_slide_content(
    Presentation("merged-table.pptx").slides[0]
)["tables"][0]
merged_origin = merged_inventory[0][0]
covered_slots = [merged_inventory[0][1], merged_inventory[1][0], merged_inventory[1][1]]
check(
    "table inventory records a merged origin and its row/column spans",
    merged_origin["is_merge_origin"]
    and not merged_origin["is_spanned"]
    and merged_origin["span_height"] == 2
    and merged_origin["span_width"] == 2,
    merged_origin,
)
check(
    "table inventory marks covered slots without repeating merged text",
    all(item["is_spanned"] and item["text"] is None for item in covered_slots)
    and sum(
        item["text"] == "Merged heading"
        for row in merged_inventory for item in row
    ) == 1,
    merged_inventory,
)

# XY scatter and bubble plots do not have category/value-series semantics.
xy_prs = Presentation()
xy_slide = xy_prs.slides.add_slide(xy_prs.slide_layouts[6])
xy_data = XyChartData()
xy_series = xy_data.add_series("XY series")
xy_series.add_data_point(1, 2)
xy_series.add_data_point(3, 4)
xy_slide.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER,
    Inches(0.5), Inches(0.5), Inches(4), Inches(2.5), xy_data,
)
bubble_data = BubbleChartData()
bubble_series = bubble_data.add_series("Bubble series")
bubble_series.add_data_point(5, 6, 7)
xy_slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE,
    Inches(0.5), Inches(3.5), Inches(4), Inches(2.5), bubble_data,
)
xy_prs.save("xy-bubble.pptx")
xy_content = extract_slide_content(Presentation("xy-bubble.pptx").slides[0])
xy_plots = [plot for chart in xy_content["charts"] for plot in chart["plots"]]
check(
    "scatter inventory emits x and y caches without category access",
    any(
        plot["kind"] == "XyPlot"
        and plot["series"][0]["x_points"] == [(0, 1.0), (1, 3.0)]
        and plot["series"][0]["y_points"] == [(0, 2.0), (1, 4.0)]
        for plot in xy_plots
    ),
    xy_plots,
)
check(
    "bubble inventory emits x, y, and bubble-size caches",
    any(
        plot["kind"] == "BubblePlot"
        and plot["series"][0]["x_points"] == [(0, 5.0)]
        and plot["series"][0]["y_points"] == [(0, 6.0)]
        and plot["series"][0]["bubble_points"] == [(0, 7.0)]
        for plot in xy_plots
    ),
    xy_plots,
)

# Missing caches and missing ptCount are unavailable, not a proven empty series.
xy_plot = next(plot for chart in Presentation("xy-bubble.pptx").slides[0].shapes
               if chart.has_chart for plot in chart.chart.plots
               if type(plot).__name__ == "XyPlot")
xy_item = list(xy_plot.series)[0]
x_source = xy_item._element.xVal
missing_cache = copy.deepcopy(x_source)
missing_cache_ref = missing_cache.find(qn("c:numRef"))
missing_cache_ref.remove(missing_cache_ref.find(qn("c:numCache")))
check("numeric cache without numCache reports unavailable", cached_numeric_points(missing_cache) is None)
missing_count = copy.deepcopy(x_source)
missing_count_cache = missing_count.find(qn("c:numRef") + "/" + qn("c:numCache"))
missing_count_cache.remove(missing_count_cache.find(qn("c:ptCount")))
check("numeric cache without ptCount reports unavailable", cached_numeric_points(missing_count) is None)
marker_source = copy.deepcopy(x_source)
marker_points = marker_source.xpath("./c:numRef/c:numCache/c:pt")
marker_points[0].find(qn("c:v")).text = "#N/A"
marker_points[1].find(qn("c:v")).text = None
check("numeric cache preserves #N/A and explicit blank markers",
      cached_numeric_points(marker_source) == [(0, "#N/A"), (1, None)],
      cached_numeric_points(marker_source))
sparse_source = copy.deepcopy(x_source)
sparse_cache = sparse_source.xpath("./c:numRef/c:numCache")[0]
sparse_cache.find(qn("c:ptCount")).set("val", "3")
sparse_cache.remove(sparse_cache.findall(qn("c:pt"))[1])
check("XY numeric cache retains only present sparse indexed points",
      cached_numeric_points(sparse_source) == [(0, 1.0)],
      cached_numeric_points(sparse_source))
check("sparse XY cache reports ptCount without densifying trailing blanks",
      cached_numeric_points(sparse_source, include_count=True)
      == {"point_count": 3, "points": [(0, 1.0)]},
      cached_numeric_points(sparse_source, include_count=True))
sparse_declared_budget = {"remaining": 3}
check("sparse numeric cache charges declared ptCount before materializing present points",
      cached_numeric_points(sparse_source, point_budget=sparse_declared_budget)
      == [(0, 1.0)] and sparse_declared_budget["remaining"] == 0,
      sparse_declared_budget)
check("category numeric cache can densify sparse positions as blanks",
      cached_numeric_points(sparse_source, fill_missing=True)
      == [(0, 1.0), (1, None), (2, None)],
      cached_numeric_points(sparse_source, fill_missing=True))
duplicate_source = copy.deepcopy(x_source)
duplicate_cache = duplicate_source.xpath("./c:numRef/c:numCache")[0]
duplicate_cache.append(copy.deepcopy(duplicate_cache.findall(qn("c:pt"))[0]))
check("duplicate numeric cache indexes fail closed",
      cached_numeric_points(duplicate_source) is None)
huge_count_source = copy.deepcopy(x_source)
huge_count_source.xpath("./c:numRef/c:numCache/c:ptCount")[0].set(
    "val", "4294967295"
)
check("oversized logical numeric counts fail before allocation",
      cached_numeric_points(huge_count_source, fill_missing=True) is None)

# Sparse XY/bubble caches reserve their declared logical counts from one shared deck
# budget before iterating or materializing the few present <c:pt> nodes.
budgeted_xy_element = copy.deepcopy(xy_item._element)
for source_name in ("xVal", "yVal"):
    numeric_cache = budgeted_xy_element.find(qn(f"c:{source_name}"))
    numeric_cache = numeric_cache.find(qn("c:numRef") + "/" + qn("c:numCache"))
    numeric_cache.find(qn("c:ptCount")).set("val", "60000")
    for point in numeric_cache.findall(qn("c:pt"))[1:]:
        numeric_cache.remove(point)
budgeted_xy_series = type("BudgetedXySeries", (), {
    "_element": budgeted_xy_element,
})()
xy_declared_budget = {"remaining": MAX_CHART_POINTS}
xy_budgeted_content = series_content(
    budgeted_xy_series, point_budget=xy_declared_budget
)
check("scatter X/Y caches charge declared sparse counts from the shared budget",
      xy_budgeted_content.get("cache_status") == "unavailable"
      and xy_declared_budget["remaining"] == 40000,
      (xy_budgeted_content, xy_declared_budget))

bubble_plot = next(
    plot for shape in Presentation("xy-bubble.pptx").slides[0].shapes
    if shape.has_chart for plot in shape.chart.plots
    if type(plot).__name__ == "BubblePlot"
)
bubble_item = list(bubble_plot.series)[0]
budgeted_bubble_element = copy.deepcopy(bubble_item._element)
for source_name in ("xVal", "yVal", "bubbleSize"):
    numeric_cache = budgeted_bubble_element.find(qn(f"c:{source_name}"))
    numeric_cache = numeric_cache.find(qn("c:numRef") + "/" + qn("c:numCache"))
    numeric_cache.find(qn("c:ptCount")).set("val", "40000")
    for point in numeric_cache.findall(qn("c:pt"))[1:]:
        numeric_cache.remove(point)
budgeted_bubble_series = type("BudgetedBubbleSeries", (), {
    "_element": budgeted_bubble_element,
})()
bubble_declared_budget = {"remaining": MAX_CHART_POINTS}
bubble_budgeted_content = series_content(
    budgeted_bubble_series, point_budget=bubble_declared_budget
)
check("bubble X/Y/size caches share one declared-count budget",
      bubble_budgeted_content.get("cache_status") == "unavailable"
      and bubble_declared_budget["remaining"] == 20000,
      (bubble_budgeted_content, bubble_declared_budget))

actual_cache = x_source.xpath("./c:numRef/c:numCache")[0]
x_source.xpath("./c:numRef")[0].remove(actual_cache)
check("XY series with an unavailable cache is explicit",
      series_content(xy_item).get("cache_status") == "unavailable", series_content(xy_item))


sp_element = slide6.shapes[-1]._element  # the textbox; layout 5 still carries a Title placeholder
group_element = parse_xml(GRP)
sp_element.getparent().replace(sp_element, group_element)
group_element.append(sp_element)

flat = list(iter_shapes(slide6.shapes))
check(
    "walker finds the shape nested in the group",
    any(getattr(sh, "text_frame", None) is not None and sh.text_frame.text == "nested member" for sh in flat),
)
check(
    "top-level shapes list hides the nested member (negative control)",
    not any(getattr(sh, "text_frame", None) is not None and sh.text_frame.text == "nested member" for sh in slide6.shapes),
)
grouped_faces = [
    run.font.name
    for shape in iter_shapes(slide6.shapes)
    if shape.has_text_frame
    for paragraph in shape.text_frame.paragraphs
    for run in paragraph.runs
    if run.font.name
]
check("font triage reaches runs nested in groups", "Grouped Face" in grouped_faces, grouped_faces)

# ---- edit.md locator: candidate collection must recurse into groups ------------
old_w, new_w = "nested member", "renamed member"


def iter_shapes_with_path(shapes, path=""):
    for shape in shapes:
        here = f"{path}/{shape.name}" if path else shape.name
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_with_path(shape.shapes, here)
        else:
            yield here, shape


def iter_text_targets(path, shape):
    if shape.has_text_frame:
        yield path, shape.text_frame
    if shape.has_table:
        for row_index, row in enumerate(shape.table.rows):
            for column_index, cell in enumerate(row.cells):
                if cell.is_spanned:
                    continue
                yield f"{path}/table[{row_index},{column_index}]", cell.text_frame


merged_edit_old, merged_edit_new = "Merged edit target", "Edited merged target"
merged_edit_prs = Presentation()
merged_edit_slide = merged_edit_prs.slides.add_slide(merged_edit_prs.slide_layouts[6])
merged_edit_table = merged_edit_slide.shapes.add_table(
    2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
).table
merged_edit_table.cell(0, 0).text = merged_edit_old
merged_edit_table.cell(0, 0).merge(merged_edit_table.cell(1, 1))
# Covered cells can retain stale text in a valid OPC package even though it is not rendered.
for merged_edit_coordinate in ((0, 1), (1, 0), (1, 1)):
    merged_edit_table.cell(*merged_edit_coordinate).text = merged_edit_old
merged_edit_prs.save("merged-edit-source.pptx")

merged_edit_reopened = Presentation("merged-edit-source.pptx")
merged_edit_shape = next(
    shape for shape in merged_edit_reopened.slides[0].shapes if shape.has_table
)
merged_edit_candidates = [
    (location, text_frame)
    for location, text_frame in iter_text_targets(merged_edit_shape.name, merged_edit_shape)
    if merged_edit_old in text_frame.text
]
check(
    "merged-table locator ignores stale text in covered merge slots",
    len(merged_edit_candidates) == 1
    and merged_edit_candidates[0][0].endswith("/table[0,0]"),
    [location for location, _ in merged_edit_candidates],
)
merged_edit_frame = merged_edit_candidates[0][1]
merged_edit_hits = [
    run
    for paragraph in merged_edit_frame.paragraphs
    for run in paragraph.runs
    if merged_edit_old in run.text
]
require(len(merged_edit_hits) == 1, "expected one run in the merged-cell edit target")
merged_edit_hits[0].text = merged_edit_hits[0].text.replace(
    merged_edit_old, merged_edit_new, 1
)
merged_edit_reopened.save("merged-edit-result.pptx")

merged_edit_result = Presentation("merged-edit-result.pptx")
merged_edit_result_table = next(
    shape.table for shape in merged_edit_result.slides[0].shapes if shape.has_table
)
merged_edit_origin = merged_edit_result_table.cell(0, 0)
merged_edit_covered = [
    merged_edit_result_table.cell(*coordinate)
    for coordinate in ((0, 1), (1, 0), (1, 1))
]
check(
    "merged-cell run edit persists without changing the merge topology",
    merged_edit_origin.text == merged_edit_new
    and merged_edit_origin.is_merge_origin
    and not merged_edit_origin.is_spanned
    and merged_edit_origin.span_width == 2
    and merged_edit_origin.span_height == 2
    and all(cell.is_spanned for cell in merged_edit_covered),
    {
        "origin": merged_edit_origin.text,
        "covered": [(cell.text, cell.is_spanned) for cell in merged_edit_covered],
    },
)


def iter_postcheck_text_frames(shapes, path=""):
    for shape in shapes:
        here = f"{path}/{shape.name}" if path else shape.name
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_postcheck_text_frames(shape.shapes, here)
            continue
        if shape.has_text_frame:
            yield here, shape.text_frame
        if shape.has_table:
            for row_index, row in enumerate(shape.table.rows):
                for column_index, cell in enumerate(row.cells):
                    if not cell.is_spanned:
                        yield f"{here}/table[{row_index},{column_index}]", cell.text_frame


group_postcheck = {
    location: frame.text for location, frame in iter_postcheck_text_frames(slide6.shapes)
}
table_postcheck = {
    location: frame.text
    for location, frame in iter_postcheck_text_frames(Presentation("input.pptx").slides[0].shapes)
}
check("mandatory postcheck reaches grouped text",
      "nested member" in group_postcheck.values(), group_postcheck)
check("mandatory postcheck reaches table-cell text",
      "old cell text" in table_postcheck.values(), table_postcheck)
check("mandatory postcheck inventories unexpected empty table cells",
      any(location.endswith("/table[0,0]") and not text.strip()
          for location, text in table_postcheck.items()), table_postcheck)


candidates = [
    (i, location, text_frame)
    for i, s in enumerate(prs6.slides)
    for p, sh in iter_shapes_with_path(s.shapes)
    for location, text_frame in iter_text_targets(p, sh)
    if old_w in text_frame.text
]
check("locator reaches text inside the group", len(candidates) == 1, [(i, p) for i, p, _ in candidates])
check("locator reports a stable nested path", "/" in candidates[0][1], candidates[0][1])
_, _, target = candidates[0]
target.paragraphs[0].runs[0].text = new_w
prs6.save("group-edited.pptx")
prs_g = Presentation("group-edited.pptx")
found = [sh for sh in iter_shapes(prs_g.slides[0].shapes)
         if getattr(sh, "text_frame", None) is not None and sh.text_frame.text == new_w]
check("group member edit persists after save", len(found) == 1)

table_candidates = [
    (i, location, text_frame)
    for i, s in enumerate(Presentation("input.pptx").slides)
    for p, sh in iter_shapes_with_path(s.shapes)
    for location, text_frame in iter_text_targets(p, sh)
    if old_cell in text_frame.text
]
check("locator reaches wording stored only in a table cell", len(table_candidates) == 1)
check("table-cell locator retains row and column",
      table_candidates[0][1].endswith("/table[0,1]"), table_candidates[0][1])

duplicate_prs = Presentation("input.pptx")
duplicate_shape = next(shape for shape in duplicate_prs.slides[0].shapes if shape.has_table)
duplicate_shape.table.cell(1, 0).text = old_cell
all_table_candidates = [
    (i, location, text_frame)
    for i, slide_item in enumerate(duplicate_prs.slides)
    for path, shape in iter_shapes_with_path(slide_item.shapes)
    for location, text_frame in iter_text_targets(path, shape)
    if old_cell in text_frame.text
]
target_location = f"{duplicate_shape.name}/table[0,1]"
selected_table_candidates = [
    candidate for candidate in all_table_candidates if candidate[1] == target_location
]
check("duplicate table text requires a location selector", len(all_table_candidates) == 2)
check("table location selector chooses one row/column",
      len(selected_table_candidates) == 1, [item[1] for item in all_table_candidates])
try:
    require(len(all_table_candidates) == 1, "target is not unique")
    optimized_duplicate_guard_rejected = False
except ValueError:
    optimized_duplicate_guard_rejected = True
check("explicit uniqueness guard rejects duplicates even under python -O",
      optimized_duplicate_guard_rejected)

split_prs = Presentation()
split_slide = split_prs.slides.add_slide(split_prs.slide_layouts[6])
split_frame = split_slide.shapes.add_textbox(
    Inches(1), Inches(1), Inches(4), Inches(1)
).text_frame
split_frame.paragraphs[0].add_run().text = "old "
split_frame.paragraphs[0].add_run().text = "wording"
split_hits = [
    run for paragraph in split_frame.paragraphs for run in paragraph.runs
    if old in run.text
]
try:
    require(len(split_hits) == 1 and split_hits[0].text.count(old) == 1,
            "target is split across runs")
    optimized_run_guard_rejected = False
except ValueError:
    optimized_run_guard_rejected = True
check("explicit run-boundary guard survives python -O", optimized_run_guard_rejected)

repeated_frame = split_slide.shapes.add_textbox(
    Inches(1), Inches(2), Inches(4), Inches(1)
).text_frame
repeated_frame.text = f"{old} / {old}"
try:
    require(repeated_frame.text.count(old) == 1,
            "target occurs more than once in the selected shape")
    optimized_repeated_guard_rejected = False
except ValueError:
    optimized_repeated_guard_rejected = True
check("explicit repeated-text guard survives python -O", optimized_repeated_guard_rejected)

# ---- analyze.md snippet: per-master, script-aware theme font resolution --------

prs7 = Presentation("input.pptx")
theme_cache = {}
DRAWINGML = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def read_theme_role(root, role):
    node = root.find(f".//a:{role}Font", DRAWINGML)
    if node is None:
        return {"latin": "", "eastAsia": "", "complexScript": "", "scripts": {}}

    def typeface(tag):
        child = node.find(f"a:{tag}", DRAWINGML)
        return "" if child is None else child.get("typeface", "")

    return {
        "latin": typeface("latin"),
        "eastAsia": typeface("ea"),
        "complexScript": typeface("cs"),
        "scripts": {
            child.get("script"): child.get("typeface", "")
            for child in node.findall("a:font", DRAWINGML) if child.get("script")
        },
    }


def theme_faces_for_slide(slide):
    master_part = slide.slide_layout.slide_master.part
    cache_key = str(master_part.partname)
    if cache_key not in theme_cache:
        theme_part = master_part.part_related_by(RT.THEME)
        root = ET.fromstring(theme_part.blob)
        theme_cache[cache_key] = {
            "major": read_theme_role(root, "major"),
            "minor": read_theme_role(root, "minor"),
        }
    return cache_key, theme_cache[cache_key]


EAST_ASIAN_SCRIPTS = {"Hans", "Hant", "Jpan", "Hang", "Bopo"}
COMPLEX_SCRIPTS = {"Arab", "Hebr", "Deva", "Beng", "Taml", "Thai"}
HAN_RANGES = (
    (0x2E80, 0x2E99), (0x2E9B, 0x2EF3), (0x2F00, 0x2FD5),
    (0x3005, 0x3005), (0x3007, 0x3007), (0x3021, 0x3029), (0x3038, 0x303B),
    (0x3400, 0x4DBF), (0x4E00, 0x9FFF),
    (0xF900, 0xFA6D), (0xFA70, 0xFAD9),
    (0x16FE2, 0x16FE3), (0x16FF0, 0x16FF6),
    (0x20000, 0x2A6DF), (0x2A700, 0x2B81D), (0x2B820, 0x2CEAD),
    (0x2CEB0, 0x2EBE0), (0x2EBF0, 0x2EE5D), (0x2F800, 0x2FA1D),
    (0x30000, 0x3134A), (0x31350, 0x33479),
)


def character_tags(character):
    codepoint = ord(character)
    if 0x3040 <= codepoint <= 0x30FF or 0x31F0 <= codepoint <= 0x31FF:
        return ["Jpan"]
    if (0x1100 <= codepoint <= 0x11FF or 0x3130 <= codepoint <= 0x318F
            or 0xA960 <= codepoint <= 0xA97F or 0xAC00 <= codepoint <= 0xD7AF
            or 0xD7B0 <= codepoint <= 0xD7FF):
        return ["Hang"]
    if 0x3100 <= codepoint <= 0x312F or 0x31A0 <= codepoint <= 0x31BF:
        return ["Bopo"]
    if any(start <= codepoint <= end for start, end in HAN_RANGES):
        return ["Hans", "Hant", "Jpan", "Hang"]
    for tag, start, end in (
        ("Cyrl", 0x0400, 0x052F), ("Hebr", 0x0590, 0x05FF),
        ("Arab", 0x0600, 0x06FF), ("Deva", 0x0900, 0x097F),
        ("Beng", 0x0980, 0x09FF), ("Taml", 0x0B80, 0x0BFF),
        ("Thai", 0x0E00, 0x0E7F),
    ):
        if start <= codepoint <= end:
            return [tag]
    return []


def script_tags(text):
    tags = []
    for character in text:
        tags.extend(character_tags(character))
    return list(dict.fromkeys(tags))


def required_slots(text):
    slots = []
    for character in text:
        tags = set(character_tags(character))
        slot = "eastAsia" if tags & EAST_ASIAN_SCRIPTS else (
            "complexScript" if tags & COMPLEX_SCRIPTS else "latin"
        )
        if slot not in slots:
            slots.append(slot)
    return slots or ["latin"]


han_examples = [
    chr(codepoint)
    for start, end in HAN_RANGES
    for codepoint in {start, end}
]
han_gap_examples = [
    chr(codepoint) for codepoint in (
        0x2E9A, 0x2EF4, 0x2FD6, 0x3006, 0x3105, 0x3131, 0x4DC0,
        0xFA6E, 0xFA6F, 0xFADA, 0x16FE1, 0x16FE4, 0x16FEF, 0x16FF7,
        0x2A6E0, 0x2B81E, 0x2CEAE, 0x2EBE1, 0x2EE5E, 0x2FA1E,
        0x3134B, 0x3347A, 0x33480,
    )
]
han_tags = {"Hans", "Hant", "Jpan", "Hang"}
supplementary_han_classified = all(
    set(character_tags(character)) == han_tags
    and required_slots(character) == ["eastAsia"]
    for character in han_examples
)
han_gaps_stay_latin = all(
    set(character_tags(character)) != han_tags
    for character in han_gap_examples
)
check("compatibility and supplementary-plane Han use the east-Asian slot",
      supplementary_han_classified)
check("Unicode 17 Han gaps and neighboring scripts are not mislabeled as Han",
      han_gaps_stay_latin)
check("Bopomofo and compatibility Jamo keep distinct East-Asian script tags",
      character_tags("ㄅ") == ["Bopo"] and required_slots("ㄅ") == ["eastAsia"]
      and character_tags("ㄱ") == ["Hang"] and required_slots("ㄱ") == ["eastAsia"])
check("Katakana phonetic extensions remain Japanese after narrowing Han ranges",
      character_tags("ㇰ") == ["Jpan"] and required_slots("ㇰ") == ["eastAsia"])
check("supplementary Han range checks remain active under optimized Python",
      __debug__ or (supplementary_han_classified and han_gaps_stay_latin))


def raw_font_slots(rpr):
    slots = {}
    if rpr is None:
        return slots
    for slot, tag in (("latin", "a:latin"), ("eastAsia", "a:ea"),
                      ("complexScript", "a:cs")):
        child = rpr.find(qn(tag))
        if child is not None and child.get("typeface"):
            slots[slot] = child.get("typeface")
    return slots


THEME_TOKENS = {
    "+mj-lt": ("major", "latin"), "+mj-ea": ("major", "eastAsia"),
    "+mj-cs": ("major", "complexScript"), "+mn-lt": ("minor", "latin"),
    "+mn-ea": ("minor", "eastAsia"), "+mn-cs": ("minor", "complexScript"),
}


def expand_theme_token(face, theme_fonts):
    role_slot = THEME_TOKENS.get(face)
    return theme_fonts[role_slot[0]][role_slot[1]] if role_slot else face


def font_candidates(run, paragraph, theme_fonts, role):
    run_slots = raw_font_slots(run._r.rPr)
    ppr = paragraph._p.pPr
    paragraph_slots = raw_font_slots(None if ppr is None else ppr.defRPr)
    tags = script_tags(run.text)
    candidates = []
    for slot in required_slots(run.text):
        fallback_role = role
        explicit_source = None
        face = run_slots.get(slot)
        if face:
            explicit_source = "run"
        else:
            face = paragraph_slots.get(slot)
            if face:
                explicit_source = "paragraph defaults"
        if face:
            resolved = expand_theme_token(face, theme_fonts)
            if resolved:
                candidates.append((slot, resolved, explicit_source))
                continue
            if face in THEME_TOKENS:
                fallback_role = THEME_TOKENS[face][0]
            else:
                continue
        role_fonts = theme_fonts[fallback_role]
        slot_tags = [
            tag for tag in tags
            if (slot == "eastAsia" and tag in EAST_ASIAN_SCRIPTS)
            or (slot == "complexScript" and tag in COMPLEX_SCRIPTS)
            or (slot == "latin" and tag not in EAST_ASIAN_SCRIPTS | COMPLEX_SCRIPTS)
        ]
        candidates.extend(
            (slot, role_fonts["scripts"][tag], f"{fallback_role} theme script {tag}")
            for tag in slot_tags if role_fonts["scripts"].get(tag)
        )
        if role_fonts.get(slot):
            candidates.append((slot, role_fonts[slot], f"{fallback_role} theme {slot}"))
    return list(dict.fromkeys(candidates))


master_name, theme_fonts = theme_faces_for_slide(prs7.slides[0])
check(
    "theme major/minor fonts resolve through the slide master relationship",
    theme_fonts["major"]["latin"] and theme_fonts["minor"]["latin"],
    (master_name, theme_fonts),
)
print("theme fonts:", ascii(theme_fonts))


class StubThemePart:
    def __init__(self, major, minor, east_asian="", complex_script="", scripts=None):
        script_nodes = "".join(
            f'<a:font script="{script}" typeface="{face}"/>'
            for script, face in (scripts or {}).items()
        )
        self.blob = (
            f'<a:theme xmlns:a="{DRAWINGML["a"]}"><a:themeElements><a:fontScheme>'
            f'<a:majorFont><a:latin typeface="{major}"/><a:ea typeface="{east_asian}"/>'
            f'<a:cs typeface="{complex_script}"/>{script_nodes}</a:majorFont>'
            f'<a:minorFont><a:latin typeface="{minor}"/><a:ea typeface="{east_asian}"/>'
            f'<a:cs typeface="{complex_script}"/>{script_nodes}</a:minorFont>'
            f'</a:fontScheme></a:themeElements></a:theme>'
        ).encode()


class StubMasterPart:
    def __init__(self, name, major, minor, **theme_options):
        self.partname = name
        self.theme_part = StubThemePart(major, minor, **theme_options)

    def part_related_by(self, relationship_type):
        assert relationship_type == RT.THEME
        return self.theme_part


def stub_slide(name, major, minor, **theme_options):
    master = type("Master", (), {"part": StubMasterPart(name, major, minor, **theme_options)})()
    layout = type("Layout", (), {"slide_master": master})()
    return type("Slide", (), {"slide_layout": layout})()


_, first_fonts = theme_faces_for_slide(stub_slide("/ppt/slideMasters/one.xml", "Head One", "Body One"))
_, second_fonts = theme_faces_for_slide(stub_slide("/ppt/slideMasters/two.xml", "Head Two", "Body Two"))
check(
    "different slide masters resolve their own theme faces",
    first_fonts["major"]["latin"] == "Head One"
    and second_fonts["major"]["latin"] == "Head Two"
    and first_fonts != second_fonts,
    (first_fonts, second_fonts),
)

_, script_fonts = theme_faces_for_slide(stub_slide(
    "/ppt/slideMasters/scripts.xml", "Latin Theme", "Latin Body",
    east_asian="East Asian Theme", complex_script="Complex Script Theme",
    scripts={
        "Hans": "Simplified Chinese Theme", "Cyrl": "Cyrillic Theme",
        "Thai": "Thai Theme",
    },
))

check(
    "unstyled runs report as inherited, not as a concrete face",
    all(run.font.name is None for sh in prs7.slides[0].shapes if sh.has_text_frame for par in sh.text_frame.paragraphs for run in par.runs),
)

font_box = prs7.slides[0].shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
font_paragraph = font_box.text_frame.paragraphs[0]
font_paragraph.font.name = "Paragraph Face"
paragraph_run = font_paragraph.add_run()
paragraph_run.text = "paragraph default"
explicit_run = font_paragraph.add_run()
explicit_run.text = "run override"
explicit_run.font.name = "Run Face"
cjk_theme_run = font_paragraph.add_run()
cjk_theme_run.text = "汉字"
thai_theme_run = font_paragraph.add_run()
thai_theme_run.text = "ไทย"
mixed_run = font_paragraph.add_run()
mixed_run.text = "A汉ก"
mixed_run.font.name = "Latin Explicit"
for tag, face in (("a:ea", "East Explicit"), ("a:cs", "Complex Explicit")):
    node = OxmlElement(tag)
    node.set("typeface", face)
    mixed_run._r.get_or_add_rPr().append(node)
token_run = font_paragraph.add_run()
token_run.text = "汉"
token_ea = OxmlElement("a:ea")
token_ea.set("typeface", "+mj-ea")
token_run._r.get_or_add_rPr().append(token_ea)

partial_box = prs7.slides[0].shapes.add_textbox(Inches(5), Inches(4), Inches(4), Inches(1))
partial_paragraph = partial_box.text_frame.paragraphs[0]
latin_only_run = partial_paragraph.add_run()
latin_only_run.text = "A汉"
latin_only_run.font.name = "Latin Only"
east_only_run = partial_paragraph.add_run()
east_only_run.text = "A汉"
east_only = OxmlElement("a:ea")
east_only.set("typeface", "East Only")
east_only_run._r.get_or_add_rPr().append(east_only)
cyrillic_cjk_run = partial_paragraph.add_run()
cyrillic_cjk_run.text = "Тест 汉"

token_fallback_fonts = {
    "major": {"latin": "Major Latin", "eastAsia": "", "complexScript": "",
              "scripts": {"Hans": "Major Hans"}},
    "minor": {"latin": "Minor Latin", "eastAsia": "Minor East", "complexScript": "",
              "scripts": {"Hans": "Minor Hans"}},
}

detected_faces = {
    run.text: font_candidates(run, font_paragraph, script_fonts, "minor")
    for run in font_paragraph.runs
}
latin_only_faces = font_candidates(latin_only_run, partial_paragraph, script_fonts, "minor")
east_only_faces = font_candidates(east_only_run, partial_paragraph, script_fonts, "minor")
cyrillic_cjk_faces = font_candidates(
    cyrillic_cjk_run, partial_paragraph, script_fonts, "minor"
)
check("font triage reports the run face and source",
      ("latin", "Run Face", "run") in detected_faces["run override"], detected_faces)
check(
    "font triage reports the paragraph face and source",
    ("latin", "Paragraph Face", "paragraph defaults")
    in detected_faces["paragraph default"],
    detected_faces,
)
check(
    "CJK inherited-font triage includes east-Asian and script-specific faces",
    {"East Asian Theme", "Simplified Chinese Theme"}
    <= {face for _, face, _ in detected_faces["汉字"]}, detected_faces["汉字"],
)
check("Thai inherited-font triage uses its script mapping",
      "Thai Theme" in {face for _, face, _ in detected_faces["ไทย"]}, detected_faces["ไทย"])
check("explicit Latin/East-Asian/complex-script slots resolve independently",
      {(slot, face) for slot, face, _ in detected_faces["A汉ก"]}
      == {("latin", "Latin Explicit"), ("eastAsia", "East Explicit"),
          ("complexScript", "Complex Explicit")}, detected_faces["A汉ก"])
check("mixed run combines direct Latin with inherited east-Asian faces",
      {("latin", "Latin Only"), ("eastAsia", "East Asian Theme")}
      <= {(slot, face) for slot, face, _ in latin_only_faces},
      latin_only_faces)
check("mixed run combines inherited Latin with a direct east-Asian face",
      {("latin", "Latin Body"), ("eastAsia", "East Only")}
      <= {(slot, face) for slot, face, _ in east_only_faces},
      east_only_faces)
check("Cyrillic plus CJK requires both Latin and east-Asian theme slots",
      {("latin", "Cyrillic Theme"), ("eastAsia", "Simplified Chinese Theme")}
      <= {(slot, face) for slot, face, _ in cyrillic_cjk_faces},
      cyrillic_cjk_faces)
check("empty +mj-ea generic face falls back to the major script mapping",
      ("eastAsia", "Major Hans", "major theme script Hans")
      in font_candidates(token_run, font_paragraph, token_fallback_fonts, "minor"),
      font_candidates(token_run, font_paragraph, token_fallback_fonts, "minor"))


# ---- analyze.md snippet: sparse cache points keep their idx ----------------------
sparse_chart_xml = (
    '<c:ser xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">'
    '<c:idx val="0"/><c:order val="0"/><c:tx><c:v>series</c:v></c:tx>'
    '<c:xVal><c:numRef><c:numCache><c:formatCode>General</c:formatCode>'
    '<c:ptCount val="6"/>'
    '<c:pt idx="0"><c:v>1</c:v></c:pt>'
    '<c:pt idx="2"><c:v>5</c:v></c:pt>'
    '<c:pt idx="5"><c:v>9</c:v></c:pt>'
    '</c:numCache></c:numRef></c:xVal></c:ser>'
)
from pptx.oxml.ns import qn as pptx_qn

def cached_numeric_points(series_element, element_name):
    return [
        (int(pt.get("idx")), pt.find(pptx_qn("c:v")).text)
        for pt in series_element.xpath(f"./c:{element_name}//c:pt")
        if pt.get("idx") is not None and pt.find(pptx_qn("c:v")) is not None
    ]

from pptx.oxml import parse_xml as pptx_parse_xml
sparse_series = pptx_parse_xml(sparse_chart_xml)
points = cached_numeric_points(sparse_series, "xVal")
check("sparse cache points keep their idx", points == [(0, "1"), (2, "5"), (5, "9")], points)
compact = [value for _, value in points]
check("compact extraction is provably lossy (negative control)",
      compact == ["1", "5", "9"] and len({idx for idx, _ in points}) == len(points))



# ---- analyze.md: script-aware run faces and table-cell triage --------------------
from pptx.oxml import parse_xml as pptx_parse_xml2
from pptx.oxml.ns import qn

def required_font_slots(text, script_tags_fn):
    tags = script_tags_fn(text)
    slots = []
    if any(ch.isascii() and ch.isalnum() for ch in text):
        slots.append("latin")
    if any(tag in ("Hans", "Hant", "Jpan", "Hang") for tag in tags):
        slots.append("eastAsia")
    if any(tag in ("Arab", "Hebr", "Deva") for tag in tags):
        slots.append("complexScript")
    return slots or ["latin"]

def explicit_run_faces(run):
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        return {}
    declared = {}
    for slot, tag in (("latin", "a:latin"), ("eastAsia", "a:ea"), ("complexScript", "a:cs")):
        node = rPr.find(qn(tag))
        if node is not None and node.get("typeface"):
            declared[slot] = node.get("typeface")
    return declared

def resolve_faces(run, text, role_fonts, script_tags_fn):
    direct = explicit_run_faces(run)
    tags = script_tags_fn(text)
    result = {}
    for slot in required_font_slots(text, script_tags_fn):
        if slot in direct:
            result[slot] = [direct[slot]]
            continue
        relevant = {
            "eastAsia": {"Hans", "Hant", "Jpan", "Hang"},
            "complexScript": {"Arab", "Hebr", "Deva"},
        }.get(slot, set())
        candidates = [role_fonts[slot]]
        candidates.extend(role_fonts["scripts"].get(tag, "") for tag in tags if tag in relevant)
        result[slot] = [face for face in dict.fromkeys(candidates) if face]
    return result

run_xml = (
    '<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:rPr lang="zh-CN"><a:latin typeface="Arial"/><a:ea typeface="SimSun"/></a:rPr>'
    '<a:t>\u6d4b\u8bd5</a:t></a:r>'
)
dual_run = pptx_parse_xml2(run_xml)
class FakeRun:
    _r = dual_run
check("a CJK run with latin+ea declared resolves to the eastAsia face",
      resolve_faces(FakeRun(), "\u6d4b\u8bd5", {
          "latin": "Theme Latin", "eastAsia": "Theme East", "complexScript": "Theme CS",
          "scripts": {},
      }, script_tags) == {"eastAsia": ["SimSun"]})
check("a mixed Latin+CJK run reports both applicable declared faces",
      resolve_faces(FakeRun(), "Q3 \u6d4b\u8bd5", {
          "latin": "Theme Latin", "eastAsia": "Theme East", "complexScript": "Theme CS",
          "scripts": {},
      }, script_tags) == {"latin": ["Arial"], "eastAsia": ["SimSun"]})

latin_only_xml = (
    '<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:rPr><a:latin typeface="Arial"/></a:rPr><a:t>Q3 \u6d4b\u8bd5</a:t></a:r>'
)
latin_only_run = pptx_parse_xml2(latin_only_xml)
class LatinOnlyRun:
    _r = latin_only_run
partial_resolution = resolve_faces(LatinOnlyRun(), "Q3 \u6d4b\u8bd5", {
    "latin": "Theme Latin", "eastAsia": "Theme East", "complexScript": "Theme CS",
    "scripts": {"Hans": "Theme Hans"},
}, script_tags)
check("a direct Latin face does not suppress inherited CJK candidates",
      partial_resolution["latin"] == ["Arial"]
      and partial_resolution["eastAsia"] == ["Theme East", "Theme Hans"],
      partial_resolution)
check("run.font.name alone would report only the Latin face (negative control)",
      dual_run.find(qn("a:rPr")).find(qn("a:latin")).get("typeface") == "Arial")

triage_deck = Presentation()
triage_slide = triage_deck.slides.add_slide(triage_deck.slide_layouts[6])
triage_shape = triage_slide.shapes.add_table(2, 2, 0, 0, 4000000, 2000000)
triage_shape.table.cell(0, 0).text_frame.text = "\u8868\u683c\u6587\u672c"
triage_deck.save("triage-table.pptx")
triage_reopened = Presentation("triage-table.pptx")
triage_table_shape = next(s for s in triage_reopened.slides[0].shapes if s.has_table)
def iter_text_frames(shapes):
    for shape in shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def unresolved_graphic_font_regions(shapes):
    for shape in shapes:
        if getattr(shape, "has_chart", False):
            yield {"shape": shape.name, "kind": "chart"}
        graphic_data = shape._element.find(".//" + qn("a:graphicData"))
        if graphic_data is not None and graphic_data.get("uri") == DIAGRAM_NS:
            yield {"shape": shape.name, "kind": "SmartArt"}


frames = list(iter_text_frames(triage_reopened.slides[0].shapes))
check("table-cell text frames are reached by the triage walker",
      any("\u8868\u683c\u6587\u672c" in f.text for f in frames), [f.text for f in frames])
check("the table graphic frame itself has no text frame (negative control)",
      not triage_table_shape.has_text_frame)

chart_shape_for_font_audit = next(
    shape for shape in Presentation("input.pptx").slides[0].shapes if shape.has_chart
)
unresolved_graphics = list(unresolved_graphic_font_regions([
    chart_shape_for_font_audit, diagram_shape, triage_table_shape,
]))
check("font triage marks both charts and SmartArt as unresolved regions",
      {item["kind"] for item in unresolved_graphics} == {"chart", "SmartArt"},
      unresolved_graphics)
check("ordinary table graphic frames are not mislabeled as chart or SmartArt",
      all(item["shape"] != triage_table_shape.name for item in unresolved_graphics),
      unresolved_graphics)
hidden_deck_unresolved_graphics = list(unresolved_graphic_font_regions(
    iter_shapes(hidden_inventory_slide.shapes)
))
check("font-region inventory excludes hidden charts and SmartArt while retaining visible siblings",
      {item["shape"] for item in hidden_deck_unresolved_graphics}
      == {visible_chart_shape.name, "Visible SmartArt"},
      hidden_deck_unresolved_graphics)
try:
    if unresolved_graphics:
        raise LookupError(f"unresolved chart/SmartArt fonts: {unresolved_graphics}")
    graphic_font_audit_rejected = False
except LookupError:
    graphic_font_audit_rejected = True
check("unresolved chart/SmartArt font regions fail closed under optimized Python",
      graphic_font_audit_rejected)

print("\n" + ("ALL PPTX FIXTURES PASSED" if not failures else f"{len(failures)} FAILURES: {failures}"))
sys.exit(0 if not failures else 1)
